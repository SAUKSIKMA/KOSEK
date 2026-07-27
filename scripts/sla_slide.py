"""
Construction de la slide de synthese "Dépassement des SLA", a partir de
l'historique stocke dans l'onglet SLA de l'Excel (cf
excel_history.write_sla_history).

Contrairement aux slides evolution/surveillance (typology_slide.py,
surveillance_slide.py), qui affichent le DERNIER mois disponible dans
l'historique avec une comparaison au mois precedent, cette slide affiche
le mois CIBLE du rapport (--year/--month), PAS le dernier mois de
l'historique -- decision du 23/06/2026 : un depassement de SLA est une
liste d'incidents nominatifs du mois en cours, pas une metrique a
comparer mois par mois, et l'absence de depassement pour le mois cible
est en soi une information utile (tableau vide), qui ne doit pas etre
masquee par un retour silencieux a un ancien mois en cas de mois cible
sans depassement (cf write_sla_history qui peut tout a fait ecrire 0
ligne pour un mois donne).

Le tableau est insere dans la 4e slide du template (deja presente, titre
"Dépassement des SLA" deja rempli, sous-titre vide a completer) -- on ne
cree pas de nouvelle slide.
"""

from openpyxl import load_workbook
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from typology_slide import (
    BRAND_GREEN, HEADER_FONT_COLOR, ROW_BAND_WHITE, ROW_BAND_GREEN_TINT,
    SUBTITLE_SHAPE_NAME, HEADER_ROW_HEIGHT,
)

SHEET_NAME = "SLA"
HEADER_FILL = BRAND_GREEN

# Badge Type SLA (colonne 1).
TYPE_SLA_STYLES = {
    "MTTA": {"fg": RGBColor(0x00, 0x3D, 0x56), "bg": RGBColor(0xDC, 0xEA, 0xF2)},
    "MTTR": {"fg": RGBColor(0xB0, 0x20, 0x20), "bg": RGBColor(0xFB, 0xE1, 0xE1)},
}

# Badge Sévérité (colonne Sévérité), reprend la palette deja utilisee
# pour les graviques de la slide surveillance.
SEVERITY_STYLES = {
    "High":          {"fg": RGBColor(0xB0, 0x20, 0x20), "bg": RGBColor(0xFB, 0xE1, 0xE1)},
    "Medium":        {"fg": RGBColor(0x8A, 0x5A, 0x00), "bg": RGBColor(0xFB, 0xEF, 0xD6)},
    "Low":           {"fg": RGBColor(0x1F, 0x7A, 0x3D), "bg": RGBColor(0xDC, 0xF0, 0xC9)},
    "Informational": {"fg": RGBColor(0x55, 0x55, 0x55), "bg": RGBColor(0xEC, 0xEC, 0xEC)},
}
_DEFAULT_BADGE = {"fg": RGBColor(0x55, 0x55, 0x55), "bg": RGBColor(0xEC, 0xEC, 0xEC)}

DATE_DISPLAY_FORMAT = "%d/%m/%Y %H:%M"


def load_sla_history(path: str) -> list:
    """
    Relit l'onglet SLA de l'Excel (cf excel_history.write_sla_history).

    Retourne une liste vide si le classeur n'a encore jamais ete alimente
    pour cet onglet (ex: --update-history jamais lance avec la slide SLA
    active) -- distinct du cas "0 depassement pour le mois cible", qui
    est gere par filter_target_month() sur une liste non vide.
    """
    wb = load_workbook(path, data_only=True)
    if SHEET_NAME not in wb.sheetnames:
        return []
    ws = wb[SHEET_NAME]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        if row[0] is None:
            continue
        rows.append({
            "Month": str(row[0]),
            "TypeSLA": row[1],
            "IncidentNumber": row[2],
            "Severity": row[3],
            "Title": row[4],
            "CreatedTime": row[5],
            "AttributionTime": row[6],
            "ClosedTime": row[7],
        })
    return rows


def filter_target_month(rows: list, year: int, month: int) -> list:
    """
    Isole les depassements du mois CIBLE du rapport (cf docstring du
    module -- PAS le dernier mois present dans l'historique).
    """
    month_str = f"{year:04d}-{month:02d}"
    return [r for r in rows if r["Month"] == month_str]


def _format_dt(value) -> str:
    if value is None or value == "":
        return ""
    if hasattr(value, "strftime"):
        return value.strftime(DATE_DISPLAY_FORMAT)
    return str(value)


def _get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_cell_text(cell, text, size_pt, bold=False, color=None, align=None, bg=None):
    cell.text = str(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    p = cell.text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    # cell.text = "" ne cree AUCUN run (verifie avec python-pptx) -- cas
    # attendu pour la colonne Clôture d'un depassement MTTA sur un
    # incident encore ouvert (cf sentinel_query.SLA_BREACHES_QUERY_TEMPLATE).
    # On applique alors juste le fond/alignement et on s'arrete la.
    if not p.runs:
        return
    run = p.runs[0]
    run.font.size = Pt(size_pt)
    run.font.name = "Arial"
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def fill_sla_slide(prs, slide, year: int, month: int, table_rows: list):
    """
    Remplit la 4e slide du template avec le tableau des incidents en
    depassement de SLA (MTTA et/ou MTTR) pour le mois cible du rapport.

    Si table_rows est vide (aucun depassement ce mois-ci), le tableau ne
    contient que la ligne d'en-tete -- c'est le comportement EXPLICITEMENT
    attendu (absence de depassement = bonne nouvelle a afficher, pas une
    slide a masquer).
    """
    month_str = f"{year:04d}-{month:02d}"

    subtitle = _get_shape_by_name(slide, SUBTITLE_SHAPE_NAME)
    if subtitle is not None:
        if table_rows:
            n = len(table_rows)
            subtitle_text = f"Mois de référence : {month_str}  ({n} dépassement{'s' if n > 1 else ''})"
        else:
            subtitle_text = f"Mois de référence : {month_str}  (aucun dépassement)"
        # Pas de mise en forme explicite sur le run -- cf typology_slide.
        # fill_evolution_slide pour le detail de la decision du 23/06/2026 :
        # on laisse le run "nu" pour heriter du format Arial Black 33pt
        # navy du placeholder (idx=10 du layout), au lieu du Arial 16pt
        # italique force precedemment.
        subtitle.text_frame.text = subtitle_text

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(1.1)

    table_top = Inches(2.6)
    table_width = slide_width - 2 * margin
    table_height = slide_height - table_top - Inches(0.6)

    n_rows = len(table_rows) + 1  # + ligne d'en-tete
    n_cols = 7

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, margin, table_top, table_width, table_height)
    table = graphic_frame.table

    col_ratios = [0.09, 0.08, 0.10, 0.32, 0.14, 0.135, 0.135]
    col_widths = [int(table_width * r) for r in col_ratios]
    col_widths[-1] = table_width - sum(col_widths[:-1])  # absorbe l'arrondi sur la derniere colonne
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Type SLA", "N°INC", "Sévérité", "Titre", "Créé le", "Attribution", "Clôture"]
    for c, header in enumerate(headers):
        _set_cell_text(table.cell(0, c), header, 14, bold=True, color=HEADER_FONT_COLOR,
                        align=PP_ALIGN.CENTER, bg=HEADER_FILL)
    table.rows[0].height = HEADER_ROW_HEIGHT

    for r, row in enumerate(table_rows, start=1):
        band = ROW_BAND_WHITE if r % 2 == 1 else ROW_BAND_GREEN_TINT

        type_style = TYPE_SLA_STYLES.get(row["TypeSLA"], _DEFAULT_BADGE)
        _set_cell_text(table.cell(r, 0), row["TypeSLA"], 13, bold=True,
                        color=type_style["fg"], align=PP_ALIGN.CENTER, bg=type_style["bg"])

        _set_cell_text(table.cell(r, 1), row["IncidentNumber"], 12, align=PP_ALIGN.CENTER, bg=band)

        sev_style = SEVERITY_STYLES.get(row["Severity"], _DEFAULT_BADGE)
        _set_cell_text(table.cell(r, 2), row["Severity"], 12, bold=True,
                        color=sev_style["fg"], align=PP_ALIGN.CENTER, bg=sev_style["bg"])

        _set_cell_text(table.cell(r, 3), row["Title"], 11, bg=band)
        _set_cell_text(table.cell(r, 4), _format_dt(row["CreatedTime"]), 11, align=PP_ALIGN.CENTER, bg=band)
        _set_cell_text(table.cell(r, 5), _format_dt(row["AttributionTime"]), 11, align=PP_ALIGN.CENTER, bg=band)
        # Clôture : vide pour un depassement MTTA dont l'incident n'est
        # pas encore clôture (cf sentinel_query.SLA_BREACHES_QUERY_TEMPLATE).
        _set_cell_text(table.cell(r, 6), _format_dt(row["ClosedTime"]), 11, align=PP_ALIGN.CENTER, bg=band)

    return slide
