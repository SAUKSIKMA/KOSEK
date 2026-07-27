"""
Génération slides COSEC depuis une requête KQL live sur Sentinel (Azure Monitor Query).
Usage : python generate_cosec.py --workspace-id <GUID>
        python generate_cosec.py --workspace-id <GUID> --tenant-id <GUID>
        python generate_cosec.py --workspace-id <GUID> --ai
        python generate_cosec.py --workspace-id <GUID> --ai --debug
"""

import json
import sys
import os
import copy
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from lxml import etree

from sentinel_query import (
    fetch_cosec_incidents, fetch_typology_history,
    fetch_severity_breakdown, fetch_classification_breakdown, fetch_resolution_times,
    fetch_sla_breaches, fetch_workspace_name,
    fetch_mitre_tactics_stats, fetch_active_rules_by_tactic,
    fetch_log_ingestion_costs,
)
from typology_normalize import aggregate_typology_rows, normalize_typology
from excel_history import write_history, write_surveillance_history, write_sla_history
from typology_slide import load_history, build_latest_with_evolution, fill_evolution_slide, move_slide_to_front
from surveillance_normalize import build_surveillance_row
from surveillance_slide import (
    load_surveillance_history,
    build_latest_with_evolution as build_latest_surveillance_evolution,
    fill_surveillance_slide,
)
from sla_normalize import build_sla_rows
from sla_slide import load_sla_history, filter_target_month, fill_sla_slide
from mitre_normalize import build_tactic_stats, build_rule_counts
from mitre_slide import fill_dispositif_surveillance_slide
from log_ingestion_normalize import build_log_ingestion_groups
from log_ingestion_slide import fill_log_ingestion_slide

TEMPLATE_PATH = "template_slide.pptx"
OUTPUT_PATH   = "COSEC_rapport.pptx"
# TEMPLATE_PATH / OUTPUT_PATH restent les valeurs par defaut (usage mono-
# client, retro-compatible) mais generate_pptx() accepte desormais
# template_path= / output_path= en parametres explicites, sur le meme
# modele que history_excel=, pour permettre l'usage multi-clients (ajoute
# le 27/07/2026 -- cf run_all_clients.py) sans dupliquer le code par client.
# Renomme le 22/06/2026 (historique_typologies.xlsx -> historique_cosec.xlsx) :
# le classeur heberge desormais 2 historiques mensuels independants, dans 2
# onglets distincts ("Typologies" et "Surveillance") -- cf excel_history.py.
HISTORY_EXCEL_PATH = "historique_cosec.xlsx"

# Bandeau "Confidentiel – COSEC" (ajoute le 28/06/2026) : prefixe fixe de
# nommage des workspaces Log Analytics, toujours tronque du nom de la
# ressource Azure pour n'en garder que le suffixe identifiant le client
# (ex: "law-prd-sentinel-emh" -> "emh" -> affiche "EMH").
WORKSPACE_NAME_PREFIX = "law-prd-sentinel-"
CONFIDENTIAL_BANNER_PREFIX = "Confidentiel"

# Mapping des colonnes entités -> libellé affiché
ENTITY_FIELDS = [
    ("Accounts",       "Accounts"),
    ("Hosts",          "Hosts"),
    ("IPs",            "IPs"),
    ("SecurityGroups", "Security Groups"),
    ("URLs",           "URLs"),
    ("Files",          "Files"),
    ("Processes",      "Processes"),
    ("CloudApps",      "Cloud Apps"),
    ("Mailboxes",      "Mailboxes"),
]


# ---------------------------------------------------------------------------
# Helpers CSV
# ---------------------------------------------------------------------------

def parse_json_array(value: str) -> list:
    """Parse un champ JSON array depuis le CSV (ex: '["val1","val2"]')."""
    if not value or value.strip() in ("", "[]"):
        return []
    try:
        parsed = json.loads(value)
        if isinstance(parsed, list):
            # Déplie les arrays imbriqués produits par make_set(tostring(...))
            result = []
            for item in parsed:
                if isinstance(item, str):
                    # Certains items sont eux-mêmes des arrays sérialisés
                    try:
                        inner = json.loads(item)
                        if isinstance(inner, list):
                            result.extend([str(v) for v in inner if v])
                        else:
                            if str(item):
                                result.append(str(item))
                    except Exception:
                        if item:
                            result.append(item)
            return result
        return [str(parsed)] if parsed else []
    except Exception:
        return [value] if value else []


def format_date(value: str) -> str:
    """
    Parse les formats de date Sentinel et retourne DD/MM/YYYY HH:MM (UTC+2 Paris).
    Formats supportés :
      - 6/16/2026, 9:34:52.250 AM  (export Sentinel CSV)
      - 2026-06-16T09:34:52.250Z   (ISO 8601)
      - 2026-06-16 09:34:52
    """
    from datetime import timezone, timedelta
    if not value:
        return ""

    paris_offset = timedelta(hours=2)  # UTC+2 (CEST, heure d'été Paris)
    value = value.strip()

    # Format export Sentinel : "6/16/2026, 9:34:52.250 AM" ou "6/16/2026, 9:34:52 AM"
    for fmt in ("%m/%d/%Y, %I:%M:%S.%f %p", "%m/%d/%Y, %I:%M:%S %p"):
        try:
            dt = datetime.strptime(value, fmt)
            dt = dt.replace(tzinfo=timezone.utc) + paris_offset
            return dt.strftime("%d/%m/%Y %H:%M")
        except ValueError:
            continue

    # Format ISO 8601 avec offset explicite (produit par datetime.isoformat()
    # sur un objet aware, ex: "2026-06-16T09:34:52.250000+00:00") — couvre
    # n'importe quel offset, pas seulement Z.
    try:
        dt = datetime.fromisoformat(value)
        from datetime import timezone as _tz
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=_tz.utc)
        dt = dt.astimezone(_tz.utc) + paris_offset
        return dt.strftime("%d/%m/%Y %H:%M")
    except ValueError:
        pass

    # Format ISO 8601 (chaînes avec Z littéral ou sans offset)
    for fmt in ("%Y-%m-%dT%H:%M:%S.%fZ", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%d %H:%M:%S"):
        try:
            dt = datetime.strptime(value, fmt)
            dt = dt.replace(tzinfo=timezone.utc) + paris_offset
            return dt.strftime("%d/%m/%Y %H:%M")
        except ValueError:
            continue

    return value


def build_entities_lines(row: dict) -> list[tuple[str, str]]:
    """
    Retourne une liste de tuples (label, valeurs_concatenees) pour les entités non vides.
    Les valeurs du même type sont regroupées sur une seule ligne.
    Ex : [("Accounts", "alice@contoso.com, bob@contoso.com"), ("Hosts", "DESKTOP-ABC")]
    """
    lines = []
    for col, label in ENTITY_FIELDS:
        values = parse_json_array(row.get(col, ""))
        cleaned = [v.strip().strip('"') for v in values if v.strip().strip('"')]
        if cleaned:
            lines.append((label, ", ".join(cleaned)))
    return lines


# ---------------------------------------------------------------------------
# Helpers PPTX / XML
# ---------------------------------------------------------------------------

def clone_slide(prs: Presentation, slide_index: int = 0):
    """Duplique la slide template et l'ajoute à la présentation.

    Au-delà de la copie du XML des shapes (spTree), il faut RECRÉER sur la
    nouvelle slide les relations de la slide source (images, hyperliens,
    SmartArt...) -- toute référence portée par un attribut r:id / r:embed /
    r:link dans le XML. Une slide créée par prs.slides.add_slide() ne
    possède au départ QUE sa relation vers le layout, pas les relations de
    la slide template ; copier le spTree tel quel (deepcopy) laisse donc
    des attributs r:embed pointant vers des rId qui n'existent pas sur la
    nouvelle slide, ou qui s'y trouvent déjà réutilisés pour une AUTRE
    relation (numérotation des rId attribuée par python-pptx au fil des
    add_slide() successifs) -- bug constaté le 29/06/2026 sur les slides
    "Focus sur incident" : à partir de la 2e slide (clonée), l'image
    centrale du visuel et les flèches qui l'entourent ne s'affichaient
    plus correctement (rId résolus vers une image différente, ou vers
    rien), alors que la toute première slide (le template lui-même,
    jamais cloné) restait correcte.

    On construit donc, pour CHAQUE relation de la slide template (sauf le
    lien vers le notesSlide -- une slide clonée ne doit pas hériter des
    notes du template), la relation correspondante sur la nouvelle slide,
    puis on mémorise la correspondance ancien rId -> nouveau rId pour
    remapper tous les attributs r:* du XML cloné.
    """
    template = prs.slides[slide_index]
    slide_layout = template.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # --- Recrée les relations (images, hyperliens, etc.) de la slide
    #     template sur la nouvelle slide, en notant ancien rId -> nouveau
    #     rId. part.relate_to() réutilise une relation existante de même
    #     type vers la même cible si elle existe déjà (ex: relation vers
    #     le layout, déjà créée par add_slide() ci-dessus) -- pas de doublon.
    NOTES_SLIDE_RELTYPE = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    )
    rid_map = {}
    for rel_id, rel in template.part.rels.items():
        if rel.reltype == NOTES_SLIDE_RELTYPE:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel_id] = new_rid

    # Copie l'arbre XML de la slide template (spTree)
    sp_tree_src = template.shapes._spTree
    sp_tree_dst = new_slide.shapes._spTree

    # Supprime les shapes auto-ajoutés par add_slide
    for el in list(sp_tree_dst):
        sp_tree_dst.remove(el)

    # Clone tous les éléments de la slide template, en remappant au passage
    # toute référence r:* (r:id, r:embed, r:link, et r:dm/r:lo/r:qs/r:cs pour
    # un éventuel SmartArt) vers le nouveau rId correspondant sur cette
    # slide -- cf rid_map construit ci-dessus.
    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for el in sp_tree_src:
        new_el = copy.deepcopy(el)
        for sub_el in new_el.iter():
            for attr, value in list(sub_el.attrib.items()):
                if attr.startswith(R_NS) and value in rid_map:
                    sub_el.set(attr, rid_map[value])
        sp_tree_dst.append(new_el)

    return new_slide


def get_shape(slide, name: str):
    """Retourne la shape par son nom."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def clear_paragraphs(tf):
    """Vide tous les paragraphes d'un text frame en gardant le premier."""
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)
    # Vide le premier paragraphe
    first_p = tf.paragraphs[0]._p
    for r in first_p.findall(qn("a:r")):
        first_p.remove(r)


def copy_rpr(source_run):
    """Clone le <a:rPr> d'un run source."""
    rpr = source_run._r.find(qn("a:rPr"))
    if rpr is not None:
        return copy.deepcopy(rpr)
    return None


def make_run(text: str, rpr_elem=None, bold: bool = False) -> etree._Element:
    """Crée un élément <a:r> avec le texte donné et un rPr optionnel."""
    r = etree.Element(qn("a:r"))
    if rpr_elem is not None:
        rpr = copy.deepcopy(rpr_elem)
        # Force ou retire le gras selon le paramètre bold
        rpr.set("b", "1" if bold else "0")
        r.append(rpr)
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return r


def make_paragraph(pPr_elem=None) -> etree._Element:
    """Crée un élément <a:p> vide avec un pPr optionnel."""
    p = etree.Element(qn("a:p"))
    if pPr_elem is not None:
        p.append(copy.deepcopy(pPr_elem))
    return p


def get_ppr(para):
    """Retourne le <a:pPr> d'un paragraphe ou None."""
    return para._p.find(qn("a:pPr"))


def get_rpr(para, run_index=0):
    """Retourne le <a:rPr> du run N d'un paragraphe ou None."""
    runs = para._p.findall(qn("a:r"))
    if run_index < len(runs):
        return runs[run_index].find(qn("a:rPr"))
    return None


def set_simple_text(tf, text: str, para_index: int = 0):
    """
    Remplace le texte d'un paragraphe existant en préservant le formatage
    du premier run.
    """
    para = tf.paragraphs[para_index]
    p = para._p
    # Récupère le rPr du premier run existant
    existing_runs = p.findall(qn("a:r"))
    rpr = None
    if existing_runs:
        rpr = existing_runs[0].find(qn("a:rPr"))
        if rpr is not None:
            rpr = copy.deepcopy(rpr)
    # Supprime tous les runs
    for r in existing_runs:
        p.remove(r)
    # Ajoute le nouveau run
    r = make_run(text, rpr)
    p.append(r)


# ---------------------------------------------------------------------------
# Remplissage de ZoneTexte 9
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Réduction automatique de la police de ZoneTexte 9 (ajouté le 28/06/2026)
# ---------------------------------------------------------------------------
#
# ZoneTexte 9 a l'autosize PowerPoint "Resize shape to fit text" (spAutoFit)
# dans le template : la FORME grandit avec le contenu, mais la SLIDE, elle,
# garde une taille fixe -- un incident avec beaucoup d'entités ou une longue
# description peut donc faire déborder la zone de texte sous le bas de la
# slide. On estime (heuristique -- aucune mesure réelle du texte rendu
# n'est possible sans passer par LibreOffice, donc estimation volontairement
# conservative côté caractères/ligne, même principe que
# typology_slide.paginate_evolution_groups) la hauteur totale nécessaire à
# la taille de police courante, et on redescend par paliers de 2pt
# (18 -> 16 -> 14 -> 12 -> 10) jusqu'à ce que ça tienne, ou jusqu'au palier
# minimal si même celui-ci ne suffit pas (mieux vaut une police un peu
# petite qu'un débordement visible hors slide).
ZONE_TEXTE_9_FONT_SIZES_PT = [18, 16, 14, 12, 10]
ZONE_TEXTE_9_BOTTOM_MARGIN = Inches(0.3)
_ZONE_TEXTE_9_CHAR_WIDTH_FACTOR = 0.58  # cf typology_slide._TYPOLOGY_CHARS_PER_LINE, même principe d'estimation


def _estimate_line_count(text: str, usable_width_in: float, font_size_pt: int) -> int:
    """Estimation conservative du nombre de lignes qu'occupera `text` une
    fois rendu, à la taille de police donnée, sur la largeur utile donnée
    -- UNIQUEMENT pour décider du palier de police de ZoneTexte 9, pas pour
    le rendu réel (laissé au wrap automatique de PowerPoint/LibreOffice)."""
    if not text:
        return 1
    avg_char_width_pt = font_size_pt * _ZONE_TEXTE_9_CHAR_WIDTH_FACTOR
    chars_per_line = max(1, int((usable_width_in * 72) / avg_char_width_pt))
    return max(1, -(-len(text) // chars_per_line))


def _estimate_zone_texte_9_height_in(logical_lines: list, usable_width_in: float, font_size_pt: int) -> float:
    """Hauteur totale estimée (en pouces) pour la liste de lignes logiques
    donnée (cf fill_zone_texte_9), à la taille de police donnée. Une ligne
    logique vide (séparateur) compte pour 1 ligne, comme dans PowerPoint."""
    line_height_in = (font_size_pt * 1.2) / 72  # interligne ~1.2x la taille de police (approximation standard)
    total_lines = sum(_estimate_line_count(text, usable_width_in, font_size_pt) for text in logical_lines)
    return total_lines * line_height_in


def fill_zone_texte_9(shape, row: dict, description_override: str = None,
                       occurrences_override: str = None, slide_height=None):
    """
    Remplit la zone texte principale (ZoneTexte 9) avec les données de l'incident.

    description_override : si fourni (texte reformulé par Claude API), remplace
    le ClassificationComment brut comme contenu du champ Description.

    occurrences_override : si fourni, remplace le champ "Occurrences" brut de
    la requête (AD.alertsCount, propre à cet incident) -- depuis le
    28/06/2026, ce champ affiche le nombre d'incidents de cette TYPOLOGIE
    sur le mois entier (même valeur que la slide "Evolution des incidents
    par typologie", cf generate_pptx/occurrence_by_typology), plus le
    compteur d'alertes de cet incident précis. Repli sur row["Occurrences"]
    si non fourni (typologie introuvable dans le lookup, ou appel sans ce
    paramètre).

    slide_height : hauteur de la slide en EMU (cf prs.slide_height) --
    nécessaire pour la réduction automatique de police si le contenu déborde
    de la slide (cf ZONE_TEXTE_9_FONT_SIZES_PT). Si non fourni, la taille de
    police maximale (18pt) est utilisée sans vérification de débordement.

    Structure cible :
      Para 0  : Typologie : <valeur>
      Para 1  : (vide)
      Para 2  : Occurrence(s) : <valeur>
      Para 3  : (vide)
      Para 4  : Log source : <valeur>
      Para 5  : (vide)
      Para 6  : Severity : <valeur>
      Para 7  : (vide)
      Para 8  : Mitre ATT&CK : <tactics> | <techniques>
      Para 9  : (vide)
      Para 10 : Entités :
      Para 11+: • <Type> : <valeur>  (une ligne par entité)
      ...
      Para N  : Catégorie de clôture : <Classification> – <ClassificationReason>
      Para N+1: (vide)
      Para N+2: Description : <ClassificationComment ou description_override>
    """
    tf = shape.text_frame
    tf.word_wrap = True

    # Sauvegarde les paragraphes originaux pour récupérer les pPr/rPr
    orig_paras = tf.paragraphs

    # --- Extraction des données ---
    title        = row.get("Title", "")
    occurrences  = occurrences_override if occurrences_override is not None else row.get("Occurrences", "")
    alert_sources = ", ".join(parse_json_array(row.get("AlertSources", "")))
    severity     = row.get("Severity", "")
    tactics      = ", ".join(parse_json_array(row.get("Tactics", "")))
    techniques   = ", ".join(parse_json_array(row.get("Techniques", "")))
    mitre        = " | ".join(filter(None, [tactics, techniques]))
    classification = row.get("Classification", "")
    reason       = row.get("ClassificationReason", "")
    categorie    = " – ".join(filter(None, [classification, reason]))

    if description_override is not None and description_override.strip():
        description = description_override.strip()
    else:
        description = row.get("ClassificationComment", "N/A") or "N/A"

    entity_lines = build_entities_lines(row)
    date_str = format_date(row.get("CreatedTime", ""))

    # --- Choix de la taille de police (cf ZONE_TEXTE_9_FONT_SIZES_PT) ---
    # Construit la liste des lignes logiques EXACTEMENT comme elles seront
    # rendues plus bas (mêmes textes combinés label+valeur), pour estimer la
    # hauteur totale à chaque palier de police et choisir le plus grand qui
    # tient dans l'espace disponible sous la zone de texte.
    logical_lines = [
        f"Date : {date_str}", "",
        f"Typologie : {title}", "",
        f"Occurrence(s) : {occurrences}", "",
        f"Log source : {alert_sources}", "",
        f"Severity : {severity}", "",
        f"Mitre ATT&CK : {mitre}", "",
        "Entités :",
    ]
    if entity_lines:
        logical_lines += [f"  • {label} : {value}" for label, value in entity_lines]
    else:
        logical_lines.append("  N/A")
    logical_lines += [
        "",
        f"Catégorie de clôture : {categorie}",
        "",
        f"Description : {description}",
    ]

    font_size_pt = ZONE_TEXTE_9_FONT_SIZES_PT[0]
    if slide_height is not None and shape.top is not None:
        usable_width_in = (shape.width - tf.margin_left - tf.margin_right) / 914400
        available_height_in = (slide_height - shape.top - ZONE_TEXTE_9_BOTTOM_MARGIN) / 914400
        font_size_pt = ZONE_TEXTE_9_FONT_SIZES_PT[-1]  # palier minimal par défaut si rien ne tient
        for candidate in ZONE_TEXTE_9_FONT_SIZES_PT:
            if _estimate_zone_texte_9_height_in(logical_lines, usable_width_in, candidate) <= available_height_in:
                font_size_pt = candidate
                break

    # --- Récupère pPr et rPr de référence depuis les paragraphes originaux ---
    ref_pPr    = get_ppr(orig_paras[0])
    ref_rPr    = get_rpr(orig_paras[0])

    # rPr labels : Arial Black + b="1" (police intrinsèquement grasse)
    rPr_bold = copy.deepcopy(ref_rPr)
    rPr_bold.set("b", "1")
    rPr_bold.set("sz", str(font_size_pt * 100))

    # rPr valeurs : Arial regular + b="0"
    # Arial Black est visuellement gras quelle que soit la valeur de b,
    # on bascule donc sur Arial pour les valeurs non-grasses.
    rPr_normal = copy.deepcopy(ref_rPr)
    rPr_normal.set("b", "0")
    rPr_normal.set("sz", str(font_size_pt * 100))
    latin_el = rPr_normal.find(qn("a:latin"))
    if latin_el is not None:
        latin_el.set("typeface", "Arial")
        for attr in ("panose", "pitchFamily", "charset"):
            if attr in latin_el.attrib:
                del latin_el.attrib[attr]

    # --- Reconstruit le txBody ---
    txBody = tf._txBody

    # Supprime tous les <a:p> existants
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    def add_label_value(label: str, value: str):
        """Ajoute un paragraphe avec label en gras + valeur en normal, puis ligne vide."""
        p = make_paragraph(ref_pPr)
        p.append(make_run(f"{label} : ", rPr_bold, bold=True))
        p.append(make_run(value, rPr_normal, bold=False))
        txBody.append(p)
        txBody.append(make_paragraph(ref_pPr))  # ligne vide

    def add_empty():
        txBody.append(make_paragraph(ref_pPr))

    # Date
    add_label_value("Date", date_str)
    # Typologie
    add_label_value("Typologie", title)
    # Occurrence(s)
    add_label_value("Occurrence(s)", occurrences)
    # Log source
    add_label_value("Log source", alert_sources)
    # Severity
    add_label_value("Severity", severity)
    # Mitre ATT&CK
    add_label_value("Mitre ATT&CK", mitre)

    # Entites (sans accent)
    p_entities = make_paragraph(ref_pPr)
    p_entities.append(make_run("Entites :", rPr_bold, bold=True))
    txBody.append(p_entities)

    if entity_lines:
        for ent_label, value in entity_lines:
            p = make_paragraph(ref_pPr)
            p.append(make_run(f"  \u2022 {ent_label} : ", rPr_bold, bold=True))
            p.append(make_run(value, rPr_normal, bold=False))
            txBody.append(p)
    else:
        p = make_paragraph(ref_pPr)
        p.append(make_run("  N/A", rPr_normal, bold=False))
        txBody.append(p)

    add_empty()

    # Categorie de cloture (sans accent)
    add_label_value("Categorie de cloture", categorie)

    # Description
    p_desc = make_paragraph(ref_pPr)
    p_desc.append(make_run("Description : ", rPr_bold, bold=True))
    p_desc.append(make_run(description, rPr_normal, bold=False))
    txBody.append(p_desc)


def fill_simple_shape(slide, shape_name: str, text: str):
    """Remplace le texte d'une shape simple (N/A -> valeur)."""
    shape = get_shape(slide, shape_name)
    if shape and shape.has_text_frame:
        set_simple_text(shape.text_frame, text)


# ---------------------------------------------------------------------------
# Bandeau "Confidentiel – COSEC - <client>" (ajoute le 28/06/2026)
# ---------------------------------------------------------------------------

def truncate_workspace_name(name: str) -> str:
    """
    Tronque le prefixe fixe de nommage des workspaces (WORKSPACE_NAME_PREFIX,
    "law-prd-sentinel-") pour n'en garder que le suffixe identifiant le
    client (ex: "law-prd-sentinel-emh" -> "EMH"), mis en MAJUSCULES pour
    l'affichage sur le bandeau "Confidentiel – COSEC".

    Si le nom ne commence pas par ce prefixe exact (cas inattendu, ex.
    convention de nommage differente sur un autre tenant), le nom complet
    est retourne tel quel (en majuscules) plutot que de tronquer a
    l'aveugle une partie potentiellement significative.

    Retourne une chaine vide si name est vide/None (cf
    sentinel_query.fetch_workspace_name, qui peut ne rien trouver).
    """
    if not name:
        return ""
    if name.lower().startswith(WORKSPACE_NAME_PREFIX):
        return name[len(WORKSPACE_NAME_PREFIX):].upper()
    return name.upper()


def _append_confidential_suffix(shapes, suffix: str) -> int:
    """Parcourt une collection de shapes et ajoute ' - <suffix>' a tout
    texte commencant par 'Confidentiel' qui ne le contient pas deja.
    Retourne le nombre de shapes modifiees."""
    updated = 0
    for shape in shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text
        if text.strip().startswith(CONFIDENTIAL_BANNER_PREFIX) and suffix not in text:
            shape.text_frame.text = f"{text.rstrip()} - {suffix}"
            updated += 1
    return updated


def update_confidential_banner(prs: Presentation, suffix: str) -> int:
    """
    Ajoute le suffixe client (ex: "EMH") au bandeau "Confidentiel – COSEC"
    affiche en haut de chaque slide, qui devient "Confidentiel – COSEC -
    EMH" -- demande du 28/06/2026.

    Le bandeau peut etre porte soit par les slide masters/layouts du pptx
    (cas le plus probable pour un element identique repete sur TOUTES les
    slides : un seul shape herite visuellement par toutes les slides sans
    etre duplique dans le XML de chaque slide), soit par un shape litteral
    present sur CHAQUE slide individuelle si le template a ete construit
    autrement. On traite les deux cas pour rester robuste a la
    construction reelle de template_slide.pptx :
      1. Slide masters ET layouts associes -- une seule modification y
         suffit alors pour TOUTES les slides, y compris celles creees
         DYNAMIQUEMENT (incidents, continuations d'evolution).
      2. Chaque slide individuelle de la presentation, AU MOMENT OU cette
         fonction est appelee -- doit donc etre appelee APRES la creation
         de toutes les slides (cf appel en toute fin de generate_pptx(),
         juste avant prs.save()), sinon les slides creees apres l'appel
         ne seraient pas couvertes.

    Idempotent (verifie que le suffixe n'est pas deja present avant
    d'ajouter) : pas de risque de double-ajout si le bandeau est trouve a
    la fois au niveau master ET slide, ou si la fonction est appelee
    plusieurs fois.

    Retourne le nombre total de shapes modifiees (a but informatif/log
    uniquement).
    """
    if not suffix:
        return 0

    total = 0
    for master in prs.slide_masters:
        total += _append_confidential_suffix(master.shapes, suffix)
        for layout in master.slide_layouts:
            total += _append_confidential_suffix(layout.shapes, suffix)

    for slide in prs.slides:
        total += _append_confidential_suffix(slide.shapes, suffix)

    return total


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def generate_pptx(workspace_id: str, year: int, month: int, tenant_id: str = None,
                   use_ai: bool = False, debug: bool = False,
                   update_history: bool = False, history_excel: str = HISTORY_EXCEL_PATH,
                   evolution_slide: bool = True, surveillance_slide: bool = True,
                   sla_slide: bool = True, dispositif_slide: bool = True,
                   log_ingestion_slide: bool = True, price_per_gb: float = 4.89,
                   template_path: str = TEMPLATE_PATH, output_path: str = OUTPUT_PATH):
    prs = Presentation(template_path)

    # --- Resolution du nom du workspace pour le bandeau "Confidentiel –
    #     COSEC - <client>" (cf update_confidential_banner, applique en
    #     toute fin de fonction). Echec tolere (warning, pas sys.exit) :
    #     c'est une information cosmetique, pas une donnee bloquante pour
    #     le reste du rapport (cf sentinel_query.fetch_workspace_name).
    workspace_name_suffix = ""
    try:
        workspace_name = fetch_workspace_name(workspace_id, tenant_id=tenant_id)
        workspace_name_suffix = truncate_workspace_name(workspace_name)
        if workspace_name_suffix:
            print(f"✅ Workspace résolu : \"{workspace_name}\" → bandeau \"{CONFIDENTIAL_BANNER_PREFIX} – COSEC - {workspace_name_suffix}\".")
        else:
            print(f"⚠ Nom du workspace introuvable pour {workspace_id} — bandeau \"Confidentiel\" non modifié.")
    except Exception as e:
        print(f"⚠ Échec de la résolution du nom du workspace : {e} — bandeau \"Confidentiel\" non modifié.")

    # --- Récupération des typologies du mois (TOUJOURS, indépendamment de
    #     --update-history) -- décision du 28/06/2026 : le champ
    #     "Occurrence(s)" des slides "Focus sur incident" affiche désormais
    #     le nombre d'incidents de cette typologie sur le mois entier (même
    #     valeur et même périmètre -- tous statuts/sévérités -- que la
    #     slide "Evolution des incidents par typologie"), et non plus le
    #     compteur d'alertes brut de l'incident lui-même (AD.alertsCount).
    #     Cette requête ne servait auparavant qu'au backfill Excel
    #     (--update-history) ; son résultat agrégé (typo_agg) est
    #     désormais réutilisé pour les DEUX usages (Excel + lookup
    #     d'occurrences) afin d'éviter une requête redondante.
    #     Échec toléré (warning, pas sys.exit) : le rapport doit pouvoir
    #     être généré même si cette requête échoue, avec un repli sur la
    #     valeur brute de chaque incident pour "Occurrence(s)".
    print(f"🔎 Récupération des typologies du mois pour {year}-{month:02d}...")
    typo_rows = []
    try:
        typo_rows = fetch_typology_history(workspace_id, year, month, tenant_id=tenant_id)
    except Exception as e:
        print(f"⚠ Échec de la requête de typologies : {e}")

    typo_agg = aggregate_typology_rows(typo_rows) if typo_rows else []
    # Lookup {typologie normalisée: nombre d'incidents} pour le mois -- cf
    # fill_zone_texte_9(occurrences_override=...). Clé = typologie APRES
    # normalisation (typology_normalize.normalize_typology), exactement
    # comme aggregate_typology_rows() la calcule : le Title brut de chaque
    # incident est donc normalisé de la même façon avant lookup (cf boucle
    # par incident plus bas), pour garantir une correspondance exacte.
    occurrence_by_typology = {r["Title"]: r["IncidentCount"] for r in typo_agg}

    if not typo_agg:
        print("⚠ Aucune donnée de typologie disponible pour ce mois — le champ \"Occurrence(s)\" des "
              "slides de détail utilisera la valeur brute de chaque incident (alertsCount), et "
              "l'historique typologies ne sera pas mis à jour même si --update-history est actif.")

    # --- Récupération des statistiques MITRE ATT&CK pour la slide
    #     "Dispositif de surveillance" (TOUJOURS, indépendamment de
    #     --update-history -- cette slide n'a PAS d'historique Excel,
    #     elle affiche directement le mois cible, comme les slides de
    #     détail par incident). mitre_rows reste à None (et non []) en
    #     cas d'échec de la requête, pour distinguer "requête en échec"
    #     (encadrés laissés à "N/A", cf mitre_slide.
    #     fill_dispositif_surveillance_slide) de "requête réussie mais 0
    #     incident ce mois-ci" (encadrés explicitement mis à "0").
    mitre_rows = None
    rules_by_tactic_raw = None
    if dispositif_slide:
        print(f"🔎 Récupération des statistiques MITRE ATT&CK pour {year}-{month:02d}...")
        try:
            mitre_rows = fetch_mitre_tactics_stats(workspace_id, year, month, tenant_id=tenant_id)
        except Exception as e:
            print(f"⚠ Échec de la requête de statistiques MITRE ATT&CK : {e} — les champs "
                  "\"Incident(s)\"/\"Dernière exécution\" de la slide \"Dispositif de surveillance\" "
                  "resteront à \"N/A\".")

        # BONUS "Nombre de règles" (cf sentinel_query.fetch_active_rules_by_tactic) :
        # échec TOLÉRÉ (warning, pas sys.exit) — information cosmétique
        # explicitement marquée optionnelle par l'utilisateur ("si non
        # récupérable, laissons tomber cette partie").
        try:
            rules_by_tactic_raw = fetch_active_rules_by_tactic(workspace_id, tenant_id=tenant_id)
        except Exception as e:
            print(f"⚠ Nombre de règles actives par tactique non récupérable ({e}) — le champ "
                  "\"Nombre de règles\" de la slide \"Dispositif de surveillance\" restera à \"N/A\".")

    # --- Récupération du coût d'ingestion des logs pour la slide "Plan
    #     de collecte" (TOUJOURS, indépendamment de --update-history --
    #     même logique que la slide MITRE ci-dessus : pas d'historique
    #     Excel pour cette slide, elle affiche directement le mois cible).
    #     price_per_gb a un defaut (4.89 €/Go, tarif contractuel fixe
    #     communique par l'utilisateur le 28/06/2026) mais reste
    #     surchargeable via --price-per-gb si ce tarif change un jour ou
    #     differe pour un autre client/contrat. None reste gere (logique
    #     "non renseigne" preservee) si la fonction est appelee directement
    #     (hors CLI) sans valeur explicite.
    log_ingestion_rows = None
    if log_ingestion_slide:
        if price_per_gb is None:
            print("⚠ --price-per-gb non fourni — la slide \"Plan de collecte\" ne sera pas renseignée "
                  "(le coût estimé par Go ingéré est requis pour cette slide).")
        else:
            print(f"🔎 Récupération du coût d'ingestion des logs pour {year}-{month:02d}...")
            try:
                log_ingestion_rows = fetch_log_ingestion_costs(
                    workspace_id, year, month, price_per_gb, tenant_id=tenant_id)
            except Exception as e:
                print(f"⚠ Échec de la requête de coût d'ingestion des logs : {e} — la slide "
                      "\"Plan de collecte\" ne sera pas renseignée.")

    # --- Mise a jour de l'historique des typologies (optionnel) ---
    # Utilise le MEME year/month que la requete d'incidents ci-dessous, par
    # construction : les deux requetes portent donc toujours sur la meme
    # periode (cf decision du 21/06/2026).
    if update_history:
        if typo_agg:
            n = write_history(history_excel, typo_agg)
            print(f"✅ Historique typologies mis à jour : {n} ligne(s) au total dans {history_excel}.")
        else:
            print("⚠ Historique typologies non modifié (aucune donnée de typologie disponible).")

        # --- Mise a jour de l'historique de surveillance (gravite / cloture /
        #     MTTA-MTTR-MTTC), MEME year/month, MEME onglet --update-history ---
        print(f"🔎 Mise à jour de l'historique surveillance pour {year}-{month:02d}...")
        try:
            severity_rows = fetch_severity_breakdown(workspace_id, year, month, tenant_id=tenant_id)
            classification_rows = fetch_classification_breakdown(workspace_id, year, month, tenant_id=tenant_id)
            resolution_times = fetch_resolution_times(workspace_id, year, month, tenant_id=tenant_id)
        except Exception as e:
            print(f"❌ Échec de la requête d'historique surveillance : {e}")
            sys.exit(1)
        if severity_rows:
            surv_row = build_surveillance_row(year, month, severity_rows, classification_rows, resolution_times)
            n = write_surveillance_history(history_excel, surv_row)
            print(f"✅ Historique surveillance mis à jour : {n} mois au total dans {history_excel}.")
        else:
            print("⚠ Aucune donnée de surveillance retournée pour ce mois — historique surveillance non modifié.")

        # --- Mise a jour de l'historique des depassements SLA (MTTA/MTTR),
        #     MEME year/month, MEME onglet --update-history, nouvel onglet
        #     "SLA" -- cf decision du 23/06/2026. write_sla_history() est
        #     appelee MEME si aucun depassement n'est trouve : un mois sans
        #     depassement doit explicitement "vider" un eventuel ancien
        #     resultat errone pour ce meme mois (cf
        #     excel_history.write_sla_history, idempotent par mois y
        #     compris pour une liste vide).
        print(f"🔎 Mise à jour de l'historique des dépassements SLA pour {year}-{month:02d}...")
        try:
            breach_rows = fetch_sla_breaches(workspace_id, year, month, tenant_id=tenant_id)
        except Exception as e:
            print(f"❌ Échec de la requête de dépassements SLA : {e}")
            sys.exit(1)
        sla_history_rows = build_sla_rows(breach_rows, year, month)
        n = write_sla_history(history_excel, year, month, sla_history_rows)
        if n:
            print(f"✅ Historique SLA mis à jour : {n} dépassement(s) pour {year}-{month:02d} dans {history_excel}.")
        else:
            print(f"✅ Historique SLA mis à jour : aucun dépassement pour {year}-{month:02d} dans {history_excel}.")

    print(f"🔎 Requête en cours sur le workspace {workspace_id} pour {year}-{month:02d}...")
    try:
        rows = fetch_cosec_incidents(workspace_id, year, month, tenant_id=tenant_id)
    except Exception as e:
        print(f"❌ Échec de la requête Sentinel : {e}")
        sys.exit(1)

    if not rows:
        print("❌ Aucun incident retourné par la requête.")
        sys.exit(1)

    print(f"✅ {len(rows)} incident(s) trouvé(s) via Azure Monitor Query.")

    # Prépare le client Claude API et l'anonymizer si l'IA est activée
    client = None
    anon = None
    if use_ai:
        from reformulate import make_client, reformulate_description
        from anonymizer import Anonymizer
        try:
            client = make_client()
        except RuntimeError as e:
            print(f"❌ {e}")
            sys.exit(1)
        anon = Anonymizer()
        print("🤖 Reformulation IA activée (Claude API)" + (" — mode debug" if debug else ""))

    # La première slide est le template — on va la remplir puis dupliquer pour les suivantes
    template_slide = prs.slides[0]

    for i, row in enumerate(rows):
        print(f"  → Incident {i+1}/{len(rows)} : {row.get('Title', '?')[:60]}")

        if i == 0:
            slide = template_slide
        else:
            slide = clone_slide(prs, slide_index=0)

        # Reformulation de la description via Claude API (si activée)
        description_override = None
        if use_ai:
            from reformulate import reformulate_description
            try:
                description_override = reformulate_description(row, anon, client, debug=debug)
            except Exception as e:
                print(f"  ⚠ Erreur reformulation IA : {e} — conservation du commentaire brut.")
                description_override = None

        # Occurrence(s) : nombre d'incidents de cette typologie sur le mois
        # entier (cf occurrence_by_typology construit plus haut), et non
        # plus le compteur d'alertes brut de cet incident -- décision du
        # 28/06/2026. Normalisation du Title identique à celle appliquée
        # côté agrégation (typology_normalize.normalize_typology), pour que
        # la clé de lookup corresponde exactement. Repli sur la valeur
        # brute (cf fill_zone_texte_9) si la typologie n'a pas de
        # correspondance (données indisponibles, ou écart de périmètre
        # improbable entre les deux requêtes).
        normalized_title = normalize_typology(row.get("Title", ""))
        occurrence_count = occurrence_by_typology.get(normalized_title)
        occurrences_override = str(occurrence_count) if occurrence_count is not None else None

        # Remplissage ZoneTexte 9 (bloc principal gauche)
        zt9 = get_shape(slide, "ZoneTexte 9")
        if zt9:
            fill_zone_texte_9(zt9, row, description_override=description_override,
                               occurrences_override=occurrences_override, slide_height=prs.slide_height)

        # Les blocs Détection / Réponse automatisée / Remédiation restent N/A pour le POC
        # (seront alimentés par Claude API dans une prochaine étape)

    # --- Slides de synthese (evolution typologie / surveillance / SLA) --
    # capturees par INDEX FIXE du template (1, 2, 3 -- ordre du template
    # original : evolution = 2e slide, surveillance = 3e slide, SLA = 4e
    # slide), AVANT tout move_slide_to_front(). Une fois qu'on commence a
    # reordonner via move_slide_to_front() (qui repose sur le slide_id,
    # pas sur la position dans prs.slides), les index de prs.slides
    # glissent et ne sont plus fiables pour de nouveaux acces par
    # position -- on capture donc les 3 references d'objets Slide UNE
    # SEULE FOIS ici (objets stables, independants de tout reordonnancement
    # ulterieur), on les remplit dans n'importe quel ordre, puis on les
    # met en avant (move_slide_to_front) dans l'ordre INVERSE de l'ordre
    # final voulu (chaque appel pousse les precedents d'un rang).
    #
    # NB (28/06/2026) : evo_slide_ref reste un index fixe UNIQUE (la 2e
    # slide du template) -- les eventuelles slides de continuation
    # "(suite)" du tableau d'evolution sont clonees DYNAMIQUEMENT par
    # fill_evolution_slide() a partir de cette reference et retournees
    # dans une liste (evo_slides, capturee plus bas), donc aucun index
    # fixe supplementaire n'est necessaire pour elles.
    #
    # mitre_slide_ref (28/06/2026) : 5e slide du template ("Dispositif de
    # surveillance"), capturee par index fixe au meme titre et au meme
    # moment que les 3 autres, pour la meme raison (stabilite avant tout
    # move_slide_to_front()).
    #
    # log_slide_ref (28/06/2026) : 6e slide du template ("Plan de
    # collecte"), meme principe.
    evo_slide_ref = prs.slides[1] if evolution_slide else None
    surv_slide_ref = prs.slides[2] if surveillance_slide else None
    sla_slide_ref = prs.slides[3] if sla_slide else None
    mitre_slide_ref = prs.slides[4] if dispositif_slide else None
    log_slide_ref = prs.slides[5] if log_ingestion_slide else None

    evo_filled = False
    evo_slides = None
    surv_filled = False
    sla_filled = False
    mitre_filled = False
    log_filled = False

    # --- Slide d'évolution par typologie (2e slide du template) ---
    if evolution_slide:
        if os.path.exists(history_excel):
            history_rows = load_history(history_excel)
            latest_month, previous_month, table_rows = build_latest_with_evolution(history_rows)
            if table_rows:
                evo_slides = fill_evolution_slide(prs, evo_slide_ref, latest_month, previous_month, table_rows)
                evo_filled = True
                comparaison = f", vs {previous_month}" if previous_month else " (premier mois suivi)"
                pagination = f", sur {len(evo_slides)} slides" if len(evo_slides) > 1 else ""
                print(f"✅ Slide d'évolution remplie pour {latest_month} "
                      f"({len(table_rows)} typologies{comparaison}{pagination}).")
            else:
                print("⚠ Historique vide — pas de slide d'évolution remplie.")
        else:
            print(f"⚠ Fichier historique {history_excel} introuvable — pas de slide d'évolution remplie. "
                  f"Utilise --update-history pour le créer.")

    # --- Slide "Etat de la surveillance" (3e slide du template) ---
    if surveillance_slide:
        if os.path.exists(history_excel):
            surveillance_rows = load_surveillance_history(history_excel)
            latest_month, previous_month, latest_row, previous_row = build_latest_surveillance_evolution(surveillance_rows)
            if latest_row:
                fill_surveillance_slide(prs, surv_slide_ref, latest_month, previous_month, latest_row, previous_row)
                surv_filled = True
                comparaison = f", vs {previous_month}" if previous_month else " (premier mois suivi)"
                print(f"✅ Slide de surveillance remplie pour {latest_month} "
                      f"({latest_row['Total']} incidents{comparaison}).")
            else:
                print("⚠ Historique surveillance vide — pas de slide de surveillance remplie.")
        else:
            print(f"⚠ Fichier historique {history_excel} introuvable — pas de slide de surveillance remplie. "
                  f"Utilise --update-history pour le créer.")

    # --- Slide "Dépassement des SLA" (4e slide du template) -- contrairement
    #     aux 2 slides ci-dessus (qui affichent le DERNIER mois de
    #     l'historique et sont sautees si cet historique est vide), celle-ci
    #     affiche le mois CIBLE du rapport (--year/--month) et est REMPLIE
    #     MEME en l'absence de depassement : un tableau vide ("aucun
    #     dépassement") est une information a part entiere, pas une raison
    #     de masquer la slide (cf sla_slide.py). Seul le cas "fichier
    #     historique introuvable" (onglet SLA jamais alimente) saute la
    #     slide, comme pour les 2 slides precedentes.
    if sla_slide:
        if os.path.exists(history_excel):
            sla_rows_all = load_sla_history(history_excel)
            sla_rows_month = filter_target_month(sla_rows_all, year, month)
            fill_sla_slide(prs, sla_slide_ref, year, month, sla_rows_month)
            sla_filled = True
            if sla_rows_month:
                print(f"✅ Slide SLA remplie pour {year}-{month:02d} "
                      f"({len(sla_rows_month)} dépassement(s)).")
            else:
                print(f"✅ Slide SLA remplie pour {year}-{month:02d} (aucun dépassement).")
        else:
            print(f"⚠ Fichier historique {history_excel} introuvable — pas de slide SLA remplie. "
                  f"Utilise --update-history pour le créer.")

    # --- Slide "Dispositif de surveillance" (5e slide du template) --
    #     contrairement aux 3 slides ci-dessus, ne lit AUCUN historique
    #     Excel : elle affiche directement les statistiques MITRE ATT&CK
    #     du mois cible, deja recuperees plus haut (mitre_rows /
    #     rules_by_tactic_raw) en meme temps que les typologies. On ne
    #     remplit donc ici que la PARTIE PPTX, a partir de ces donnees
    #     deja en memoire.
    if dispositif_slide:
        tactic_stats = build_tactic_stats(mitre_rows) if mitre_rows is not None else None
        rules_by_tactic = build_rule_counts(rules_by_tactic_raw) if rules_by_tactic_raw is not None else None

        if tactic_stats is None and rules_by_tactic is None:
            print("⚠ Aucune donnée disponible pour la slide \"Dispositif de surveillance\" — "
                  "tous les encadrés restent à \"N/A\".")
        else:
            n_filled, unmatched = fill_dispositif_surveillance_slide(
                mitre_slide_ref, tactic_stats=tactic_stats, rules_by_tactic=rules_by_tactic)
            mitre_filled = True
            total_incidents = sum(s["incident_count"] for s in tactic_stats.values()) if tactic_stats else 0
            print(f"✅ Slide \"Dispositif de surveillance\" remplie pour {year}-{month:02d} "
                  f"({n_filled} encadré(s) de tactique, {total_incidents} incident(s) au total).")
            if unmatched:
                print(f"⚠ Tactique(s) sans encadré correspondant dans le template, ignorée(s) : "
                      f"{', '.join(unmatched)}.")

    # --- Slide "Plan de collecte" (6e slide du template) -- comme
    #     Dispositif de surveillance ci-dessus, ne lit AUCUN historique
    #     Excel : affiche directement le cout d'ingestion des logs du
    #     mois cible, deja recupere plus haut (log_ingestion_rows).
    if log_ingestion_slide:
        if log_ingestion_rows is not None:
            groups = build_log_ingestion_groups(log_ingestion_rows)
            fill_log_ingestion_slide(prs, log_slide_ref, year, month, groups)
            log_filled = True
            if groups:
                total_size = sum(g["size_gb"] for g in groups)
                total_cost = sum(g["cost"] for g in groups)
                print(f"✅ Slide \"Plan de collecte\" remplie pour {year}-{month:02d} "
                      f"({len(groups)} catégorie(s) de logs, {total_size:.2f} Go, "
                      f"{total_cost:.2f} € estimés).")
            else:
                print(f"✅ Slide \"Plan de collecte\" remplie pour {year}-{month:02d} "
                      "(aucune ingestion facturable trouvée).")
        else:
            print("⚠ Slide \"Plan de collecte\" non renseignée (cf avertissement ci-dessus).")

    # --- Mise en avant des slides de synthese deja remplies, dans l'ordre
    #     INVERSE de l'ordre final voulu (cf decision du 22/06/2026 pour
    #     Evolution/Surveillance, etendue le 23/06/2026 a SLA, le
    #     28/06/2026 a Dispositif de surveillance, puis a Plan de
    #     collecte) : appeler move_slide_to_front() sur Plan de collecte,
    #     puis Dispositif, puis SLA, puis Evolution, puis Surveillance
    #     donne l'ordre final Surveillance, Evolution, SLA, Dispositif de
    #     surveillance, Plan de collecte, [details incidents...] --
    #     chaque appel pousse les precedents d'un rang.
    #
    #     Plan de collecte est poussee EN PREMIER (donc la plus proche
    #     des slides de detail dans l'ordre final) : c'est la derniere
    #     des 5 slides de synthese dans l'ordre du template original
    #     (apres Evolution/Surveillance/SLA/Dispositif), meme logique que
    #     pour Dispositif de surveillance lors de son ajout.
    #
    #     Cas Evolution multi-slides (decision du 28/06/2026, cf
    #     typology_slide.fill_evolution_slide) : evo_slides contient
    #     potentiellement PLUSIEURS slides (slide principale + slides de
    #     continuation "(suite)"). On les met en avant dans l'ordre
    #     INVERSE de leur ordre final voulu (cf docstring de
    #     fill_evolution_slide) -- meme logique que pour les 3 sections,
    #     simplement appliquee a une liste plutot qu'a une slide unique.
    if log_filled:
        move_slide_to_front(prs, log_slide_ref)
    if mitre_filled:
        move_slide_to_front(prs, mitre_slide_ref)
    if sla_filled:
        move_slide_to_front(prs, sla_slide_ref)
    if evo_filled:
        for evo_slide in reversed(evo_slides):
            move_slide_to_front(prs, evo_slide)
    if surv_filled:
        move_slide_to_front(prs, surv_slide_ref)

    # --- Bandeau "Confidentiel – COSEC - <client>" -- appele en TOUT
    #     DERNIER, une fois que toutes les slides existent (incidents,
    #     synthese, continuations d'evolution) : cf update_confidential_
    #     banner(), qui doit voir l'integralite des slides de prs.slides.
    if workspace_name_suffix:
        n_banner = update_confidential_banner(prs, workspace_name_suffix)
        print(f"✅ Bandeau \"Confidentiel – COSEC\" mis à jour avec le suffixe \"{workspace_name_suffix}\" "
              f"({n_banner} zone(s) de texte modifiée(s)).")

    prs.save(output_path)
    print(f"\n✅ Fichier généré : {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Génère les slides COSEC depuis une requête KQL live sur Sentinel.")
    parser.add_argument("--workspace-id", required=True,
                         help="GUID du workspace Log Analytics (Overview > Workspace ID)")
    parser.add_argument("--tenant-id", default=None,
                         help="GUID du tenant cible (optionnel — force la résolution si accès multi-tenant via Lighthouse)")
    parser.add_argument("--year", type=int, required=True,
                         help="Année du mois cible (ex: 2026) — utilisée pour les slides de détail ET la slide d'évolution")
    parser.add_argument("--month", type=int, required=True,
                         help="Mois cible 1-12 — utilisée pour les slides de détail ET la slide d'évolution")
    parser.add_argument("--ai", action="store_true",
                         help="Active la reformulation de la description via Claude API (anonymisée)")
    parser.add_argument("--debug", action="store_true",
                         help="Mode debug : demande une validation humaine avant chaque envoi à Claude API")
    parser.add_argument("--update-history", action="store_true",
                         help="Récupère l'historique des typologies, de surveillance ET des dépassements SLA "
                              "pour --year/--month et les intègre dans l'Excel avant de générer le pptx")
    parser.add_argument("--history-excel", default=HISTORY_EXCEL_PATH,
                         help=f"Chemin du fichier Excel historique, onglets Typologies + Surveillance + SLA "
                              f"(défaut : {HISTORY_EXCEL_PATH})")
    parser.add_argument("--no-evolution-slide", action="store_true",
                         help="Ne pas ajouter la slide d'évolution par typologie")
    parser.add_argument("--no-surveillance-slide", action="store_true",
                         help="Ne pas ajouter la slide \"Etat de la surveillance\" (gravité/clôture/MTTA-MTTR-MTTC)")
    parser.add_argument("--no-sla-slide", action="store_true",
                         help="Ne pas ajouter la slide \"Dépassement des SLA\" (incidents en dépassement MTTA/MTTR)")
    parser.add_argument("--no-dispositif-slide", action="store_true",
                         help="Ne pas ajouter la slide \"Dispositif de surveillance\" "
                              "(couverture par tactique MITRE ATT&CK)")
    parser.add_argument("--no-log-ingestion-slide", action="store_true",
                         help="Ne pas ajouter la slide \"Plan de collecte\" "
                              "(coût de l'ingestion des logs par catégorie)")
    parser.add_argument("--price-per-gb", type=float, default=4.89,
                         help="Prix en €/Go ingéré, pour le calcul du coût estimé de la slide "
                              "\"Plan de collecte\" (défaut : 4.89 €/Go) — passer une autre valeur "
                              "pour l'écraser ponctuellement")
    parser.add_argument("--template-path", default=TEMPLATE_PATH,
                         help=f"Chemin du template pptx (défaut : {TEMPLATE_PATH}) — utile pour "
                              f"pointer vers un template partagé situé ailleurs que le dossier courant")
    parser.add_argument("--output", default=OUTPUT_PATH,
                         help=f"Chemin du fichier pptx de sortie (défaut : {OUTPUT_PATH}) — utile "
                              f"pour générer des rapports nommés/rangés par client")
    args = parser.parse_args()

    if args.debug and not args.ai:
        print("⚠ --debug n'a d'effet qu'avec --ai. Ajout automatique de --ai.")
        args.ai = True

    generate_pptx(args.workspace_id, args.year, args.month, tenant_id=args.tenant_id,
                  use_ai=args.ai, debug=args.debug,
                  update_history=args.update_history, history_excel=args.history_excel,
                  evolution_slide=not args.no_evolution_slide,
                  surveillance_slide=not args.no_surveillance_slide,
                  sla_slide=not args.no_sla_slide,
                  dispositif_slide=not args.no_dispositif_slide,
                  log_ingestion_slide=not args.no_log_ingestion_slide,
                  price_per_gb=args.price_per_gb,
                  template_path=args.template_path,
                  output_path=args.output)
