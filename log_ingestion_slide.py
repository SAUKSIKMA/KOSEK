"""
Construction de la slide "Plan de collecte" (6e slide du template,
ajoutee le 28/06/2026), qui reproduit -- sous forme statique, adaptee a
une slide PowerPoint -- la visualisation du workbook source fournie par
l'utilisateur (capture d'ecran) : un tableau arborescent categorie/table,
avec une barre proportionnelle a la taille ingeree et un degrade de
couleur vert/orange/rouge (cf log_ingestion_normalize.heat_fraction).

Le tableau lui-meme (texte, en-tete, bandes alternees) reprend le style
deja etabli dans typology_slide.py/sla_slide.py (couleurs de marque,
bandes blanc/vert pale) plutot que le gris neutre du portail Azure, pour
rester coherent avec le reste du rapport -- seule la partie qui PORTE une
information (barre + degrade de couleur) est reprise du workbook.

La colonne "Taille de la table" reste VIDE cote tableau natif pptx (pas
de texte dans la cellule) : la barre ET son libelle numerique sont des
shapes separees, positionnees en valeurs absolues par-dessus le tableau,
car python-pptx ne permet pas de dessiner une barre de largeur partielle
A L'INTERIEUR d'une cellule -- seul un remplissage plein-largeur est
possible nativement. Le libelle est colle juste apres la fin de la barre
(comme dans la capture d'ecran source), pas a une position fixe.

Ligne d'en-tete de categorie vs ligne "Sous-total" (decision du
29/06/2026) : la ligne d'en-tete de chaque categorie (ex: "Azure Active
Directory (2)") n'affiche QUE ce libelle, sur une cellule fusionnee sur
les 3 colonnes, fond GROUP_HEADER_FILL (#4F81BD, bleu) -- ni barre, ni
cout. La barre + le cout CUMULES de la categorie sont portes par une
ligne "Sous-total" dediee, ajoutee APRES les tables de cette categorie,
sur fond clair SUBTOTAL_FILL (#D9E6EC) : le degrade vert/orange/rouge de
la barre, et le libelle numerique en texte BRAND_NAVY, etaient peu
lisibles sur le fond bleu fonce de l'en-tete (cf retour utilisateur).
"""

import copy

from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from typology_slide import (
    BRAND_GREEN, BRAND_NAVY, ROW_BAND_WHITE, ROW_BAND_GREEN_TINT,
    HEADER_FONT_COLOR, HEADER_ROW_HEIGHT, SUBTOTAL_FILL, SUBTOTAL_FONT_COLOR,
    GROUP_HEADER_FILL, GROUP_HEADER_FONT_COLOR,
)
from log_ingestion_normalize import format_size, format_cost, heat_fraction, select_groups_for_display

CAPTION_SHAPE_NAME = "ZoneTexte 8"
HEADER_FILL = BRAND_GREEN

# Degrade vert (faible) -> orange (moyen) -> rouge (eleve), applique sur
# heat_fraction() (0.0 a 1.0) -- 2 segments lineaires (0-0.5 et 0.5-1.0).
_HEAT_GREEN = RGBColor(0x1F, 0x7A, 0x3D)
_HEAT_ORANGE = RGBColor(0xE8, 0x9A, 0x1C)
_HEAT_RED = RGBColor(0xB0, 0x20, 0x20)

DATA_ROW_HEIGHT = Inches(0.46)
BAR_ZONE_RATIO = 0.62  # fraction de la largeur de la colonne "Taille" reservee a la barre la plus longue
BAR_HEIGHT_RATIO = 0.5  # fraction de la hauteur de ligne occupee par la barre
BAR_LABEL_GAP = Inches(0.08)

_MONTH_NAMES_FR = {
    1: "janvier", 2: "février", 3: "mars", 4: "avril", 5: "mai", 6: "juin",
    7: "juillet", 8: "août", 9: "septembre", 10: "octobre", 11: "novembre", 12: "décembre",
}

_CAPTION_TEMPLATE = ("Ici, nous pouvons retrouver la répartition des ingestions facturables "
                      "par catégorie de journaux (log) au cours du mois de {month_name}")


def _get_shape_by_name(slide, name):
    """Retourne la shape de la slide portant le nom donne, ou None si aucune
    shape ne correspond (cf typology_slide._get_shape_by_name, duplique ici)."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_textbox_text(shape, text: str):
    """Remplace le texte d'un textbox simple en preservant la mise en
    forme du run existant. ZoneTexte 8 n'est PAS un placeholder (qui
    heriterait du layout) : une simple affectation de text_frame.text
    perdrait sa mise en forme -- meme principe que generate_cosec.
    set_simple_text(), duplique ici pour eviter une dependance croisee
    vers generate_cosec.py (cf mitre_normalize._parse_iso/sla_normalize.
    _parse_iso pour le meme choix de duplication ailleurs dans le
    projet)."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    existing_runs = p._p.findall(qn("a:r"))
    rpr_xml = None
    if existing_runs:
        rpr_el = existing_runs[0].find(qn("a:rPr"))
        if rpr_el is not None:
            rpr_xml = copy.deepcopy(rpr_el)
    for r in existing_runs:
        p._p.remove(r)
    new_r = etree.SubElement(p._p, qn("a:r"))
    if rpr_xml is not None:
        new_r.append(rpr_xml)
    t = etree.SubElement(new_r, qn("a:t"))
    t.text = text


def _heat_color(fraction: float) -> RGBColor:
    """Interpole la couleur de barre sur le degrade vert/orange/rouge
    (cf docstring du module) pour une fraction (0.0 a 1.0) donnee."""
    if fraction <= 0.5:
        t = fraction / 0.5
        lo, hi = _HEAT_GREEN, _HEAT_ORANGE
    else:
        t = (fraction - 0.5) / 0.5
        lo, hi = _HEAT_ORANGE, _HEAT_RED
    r = int(lo[0] + (hi[0] - lo[0]) * t)
    g = int(lo[1] + (hi[1] - lo[1]) * t)
    b = int(lo[2] + (hi[2] - lo[2]) * t)
    return RGBColor(r, g, b)


def _set_cell_text(cell, text, size_pt, bold=False, color=None, align=None, bg=None):
    """Remplit une cellule de tableau pptx (texte, fond, alignement, police/
    taille/gras/couleur du run) -- cf typology_slide._set_cell_text, duplique
    ici avec une garde sur runs vides (utile pour les colonnes sans texte,
    ex: colonne "Taille" remplie uniquement par la barre de heatmap)."""
    cell.text = str(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    p = cell.text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    if not p.runs:
        return
    run = p.runs[0]
    run.font.size = Pt(size_pt)
    run.font.name = "Arial"
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _max_rows_for_height(available_height) -> int:
    """Nombre maximal de lignes de DONNEES (hors en-tete) qui tiennent
    dans la hauteur disponible, a DATA_ROW_HEIGHT fixe."""
    return max(1, int(available_height // DATA_ROW_HEIGHT))


def fill_log_ingestion_slide(prs, slide, year: int, month: int, groups: list):
    """
    Remplit la 6e slide du template ("Plan de collecte") avec :
      - la legende (ZoneTexte 8), mois cible substitue dynamiquement ;
      - un tableau arborescent categorie/table avec barre proportionnelle
        et degrade de couleur (cf docstring du module).

    groups : sortie de log_ingestion_normalize.build_log_ingestion_groups().
    Si vide (aucune donnee d'ingestion facturable trouvee pour le mois),
    le tableau n'est PAS cree -- seule la legende est mise a jour, avec
    une mention explicite d'absence de donnees plutot qu'un tableau vide
    (a la difference de la slide SLA, ou un tableau vide est lui-meme une
    information utile -- ici, une absence totale de donnees signale plus
    vraisemblablement un probleme de requete/perimetre qu'un mois sans
    ingestion, donc une mention textuelle est plus appropriee qu'un
    tableau a une seule ligne d'en-tete).
    """
    caption = _get_shape_by_name(slide, CAPTION_SHAPE_NAME)
    month_name = _MONTH_NAMES_FR.get(month, str(month))
    if caption is not None:
        _set_textbox_text(caption, _CAPTION_TEMPLATE.format(month_name=month_name))

    if not groups:
        if caption is not None:
            _set_textbox_text(caption, _CAPTION_TEMPLATE.format(month_name=month_name)
                               + " — aucune donnée d'ingestion facturable trouvée sur cette période.")
        return slide

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(1.6)

    table_top = Inches(3.0)
    table_width = slide_width - 2 * margin
    available_height = slide_height - table_top - Inches(0.6)

    # -1 : reserve une ligne pour le total general (cf rendu en fin de
    # fonction), qui s'ajoute TOUJOURS apres les categories/le resume,
    # contrairement a ces dernieres qui sont elles soumises au budget.
    max_rows = max(1, _max_rows_for_height(available_height - HEADER_ROW_HEIGHT) - 1)
    visible_groups, summary_row = select_groups_for_display(groups, max_rows)

    n_data_rows = sum(2 + len(g["tables"]) for g in visible_groups) + (1 if summary_row else 0) + 1  # +1 par groupe : ligne "Sous-total" ; +1 final : ligne "Total"
    n_rows = 1 + n_data_rows
    n_cols = 3

    table_height = HEADER_ROW_HEIGHT + n_data_rows * DATA_ROW_HEIGHT
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, margin, table_top, table_width, table_height)
    table = graphic_frame.table

    col_ratios = [0.42, 0.36, 0.22]
    col_widths = [int(table_width * r) for r in col_ratios]
    col_widths[-1] = table_width - sum(col_widths[:-1])
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Type de journal", "Taille de la table", "Coût estimé"]
    for c, header in enumerate(headers):
        _set_cell_text(table.cell(0, c), header, 14, bold=True, color=HEADER_FONT_COLOR,
                        align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER, bg=HEADER_FILL)
    table.rows[0].height = HEADER_ROW_HEIGHT

    # Bornes globales (categories ET tables confondues) de l'echelle
    # barre/couleur -- cf log_ingestion_normalize.heat_fraction, qui a
    # besoin du min ET du max du jeu de donnees AFFICHE (pas d'une
    # constante absolue) pour etaler les couleurs sur toute l'amplitude
    # reellement presente ce mois-ci. Inclut le total des categories
    # masquees (summary_row) pour rester juste si cette ligne de synthese
    # s'avere etre la plus grosse ou la plus petite valeur (cas peu
    # probable mais possible).
    all_values = [g["size_gb"] for g in visible_groups]
    for g in visible_groups:
        all_values += [t["size_gb"] for t in g["tables"]]
    if summary_row:
        all_values.append(summary_row["size_gb"])
    max_size_gb = max(all_values) if all_values else 0.0
    min_size_gb = min(all_values) if all_values else 0.0

    bar_zone_left = margin + col_widths[0]
    bar_zone_width = int(col_widths[1] * BAR_ZONE_RATIO)
    bar_height = int(DATA_ROW_HEIGHT * BAR_HEIGHT_RATIO)

    def _draw_bar_and_label(row_top, size_gb):
        fraction = heat_fraction(size_gb, min_size_gb, max_size_gb)
        bar_width = max(Emu(1), int(bar_zone_width * fraction))
        bar_top = row_top + (DATA_ROW_HEIGHT - bar_height) // 2

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_zone_left, bar_top, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = _heat_color(fraction)
        bar.line.fill.background()
        bar.shadow.inherit = False

        label_left = bar_zone_left + bar_width + BAR_LABEL_GAP
        label_width = col_widths[1] - bar_width - BAR_LABEL_GAP
        if label_width > 0:
            label_box = slide.shapes.add_textbox(label_left, row_top, label_width, DATA_ROW_HEIGHT)
            ltf = label_box.text_frame
            ltf.word_wrap = False
            ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ltf.margin_left = 0
            ltf.margin_right = 0
            lp = ltf.paragraphs[0]
            lp.alignment = PP_ALIGN.LEFT
            lr = lp.add_run()
            lr.text = format_size(size_gb)
            lr.font.name = "Arial"
            lr.font.size = Pt(12)
            lr.font.color.rgb = BRAND_NAVY

    row_top = table_top + HEADER_ROW_HEIGHT
    r = 1

    for g in visible_groups:
        # Demande du 29/06/2026 : la ligne d'en-tete de categorie n'affiche
        # plus QUE le libelle (cellule fusionnee sur les 3 colonnes, fond
        # GROUP_HEADER_FILL -- meme principe que les en-tetes de groupe de
        # la slide "Evolution des incidents par typologie") -- ni barre, ni
        # cout. Sur le fond bleu fonce, le degrade vert/orange/rouge de la
        # barre ET le libelle numerique (texte BRAND_NAVY) etaient peu
        # lisibles (cf retour utilisateur). La valeur cumulee de la
        # categorie (barre + cout) est desormais portee par une ligne
        # "Sous-total" dediee, ajoutee APRES les tables de cette categorie,
        # sur fond clair SUBTOTAL_FILL (#D9E6EC) -- bien plus contraste.
        table.cell(r, 0).merge(table.cell(r, n_cols - 1))
        _set_cell_text(table.cell(r, 0), f"{g['category']} ({len(g['tables'])})", 13, bold=True,
                        color=GROUP_HEADER_FONT_COLOR, align=PP_ALIGN.LEFT, bg=GROUP_HEADER_FILL)
        table.rows[r].height = DATA_ROW_HEIGHT
        row_top += DATA_ROW_HEIGHT
        r += 1

        for t in g["tables"]:
            band = ROW_BAND_WHITE if (r % 2 == 1) else ROW_BAND_GREEN_TINT
            cell0 = table.cell(r, 0)
            _set_cell_text(cell0, t["name"], 12, align=PP_ALIGN.LEFT, bg=band)
            cell0.text_frame.margin_left = Inches(0.35)  # indentation visuelle des tables sous leur categorie
            table.cell(r, 1).fill.solid()
            table.cell(r, 1).fill.fore_color.rgb = band
            _set_cell_text(table.cell(r, 2), format_cost(t["cost"]), 12, align=PP_ALIGN.RIGHT, bg=band)
            table.rows[r].height = DATA_ROW_HEIGHT
            _draw_bar_and_label(row_top, t["size_gb"])
            row_top += DATA_ROW_HEIGHT
            r += 1

        # --- Ligne "Sous-total" de la categorie (cf commentaire ci-dessus) :
        #     porte desormais la barre + le libelle de taille cumules, qui
        #     etaient auparavant sur la ligne d'en-tete -- fond SUBTOTAL_FILL,
        #     bien plus contraste pour le degrade et le texte BRAND_NAVY.
        _set_cell_text(table.cell(r, 0), "Sous-total", 13, bold=True,
                        color=SUBTOTAL_FONT_COLOR, align=PP_ALIGN.LEFT, bg=SUBTOTAL_FILL)
        table.cell(r, 1).fill.solid()
        table.cell(r, 1).fill.fore_color.rgb = SUBTOTAL_FILL
        _set_cell_text(table.cell(r, 2), format_cost(g["cost"]), 13, bold=True,
                        color=SUBTOTAL_FONT_COLOR, align=PP_ALIGN.RIGHT, bg=SUBTOTAL_FILL)
        table.rows[r].height = DATA_ROW_HEIGHT
        _draw_bar_and_label(row_top, g["size_gb"])
        row_top += DATA_ROW_HEIGHT
        r += 1

    if summary_row:
        band = RGBColor(0xEC, 0xEC, 0xEC)
        muted = RGBColor(0x6B, 0x6B, 0x6B)
        _set_cell_text(table.cell(r, 0), summary_row["category"], 12, bold=False,
                        color=muted, align=PP_ALIGN.LEFT, bg=band)
        table.cell(r, 1).fill.solid()
        table.cell(r, 1).fill.fore_color.rgb = band
        _set_cell_text(table.cell(r, 2), format_cost(summary_row["cost"]), 12,
                        color=muted, align=PP_ALIGN.RIGHT, bg=band)
        table.rows[r].height = DATA_ROW_HEIGHT
        _draw_bar_and_label(row_top, summary_row["size_gb"])
        row_top += DATA_ROW_HEIGHT
        r += 1

    # --- Ligne "Total" : porte sur l'INTEGRALITE de groups (pas
    #     seulement visible_groups + summary_row, meme si les deux sont
    #     normalement equivalents en somme) -- calcul direct et explicite
    #     depuis la donnee source, plus robuste a toute evolution future
    #     de select_groups_for_display(). Pas de barre (cf docstring du
    #     module : la barre/degrade compare les lignes ENTRE ELLES, le
    #     total general n'a pas vocation a etre compare au reste sur la
    #     meme echelle) -- texte simple, mis en avant par un fond marque
    #     distinct (bleu GROUP_HEADER_FILL, demande du 29/06/2026) plutot
    #     que par une barre.
    total_size_gb = sum(g["size_gb"] for g in groups)
    total_cost = sum(g["cost"] for g in groups)
    _set_cell_text(table.cell(r, 0), "Total", 14, bold=True,
                    color=GROUP_HEADER_FONT_COLOR, align=PP_ALIGN.LEFT, bg=GROUP_HEADER_FILL)
    _set_cell_text(table.cell(r, 1), format_size(total_size_gb), 14, bold=True,
                    color=GROUP_HEADER_FONT_COLOR, align=PP_ALIGN.CENTER, bg=GROUP_HEADER_FILL)
    _set_cell_text(table.cell(r, 2), format_cost(total_cost), 14, bold=True,
                    color=GROUP_HEADER_FONT_COLOR, align=PP_ALIGN.RIGHT, bg=GROUP_HEADER_FILL)
    table.rows[r].height = DATA_ROW_HEIGHT

    return slide
