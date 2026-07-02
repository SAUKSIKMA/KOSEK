"""
Construction de la slide de synthese "Etat de la surveillance" (3e slide
du template), a partir de l'historique stocke dans l'onglet Surveillance
de l'Excel (cf excel_history.write_surveillance_history).

Comme la slide d'evolution par typologie (cf typology_slide.py), cette
slide affiche le DERNIER mois disponible dans l'historique avec une
evolution par rapport au mois precedent. Le tableau est insere dans la
3e slide du template (deja presente, titre "Les Incidents du mois" deja
rempli, sous-titre vide a completer) -- on ne cree pas de nouvelle slide.

Contenu, conformement a la demande du 22/06/2026 :
  - repartition par gravite (liste + donut)
  - MTTA / MTTR / MTTC moyens (3 cartes KPI)
  - repartition par categorie de cloture (liste + donut)
"""

from openpyxl import load_workbook
from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

from typology_slide import (
    BRAND_GREEN, BRAND_NAVY, EVOLUTION_STYLES, SUBTITLE_SHAPE_NAME,
)
from surveillance_normalize import SEVERITY_LABELS, SEVERITY_ORDER

SHEET_NAME = "Surveillance"

PANEL_BG = RGBColor(0xF2, 0xF2, 0xF2)
TEXT_NAVY = BRAND_NAVY
TEXT_MUTED = RGBColor(0x6B, 0x6B, 0x6B)

# Couleurs des 4 tranches de gravite (donut 1), reprises de l'identite
# visuelle existante (vert marque + degrade neutre/alerte).
SEVERITY_COLORS = {
    "High": RGBColor(0xB0, 0x20, 0x20),
    "Medium": RGBColor(0xE8, 0x9A, 0x1C),
    "Low": BRAND_GREEN,
    "Informational": RGBColor(0xB7, 0xB7, 0xB7),
}

# Couleurs des 4 categories de cloture (donut 2).
CLASSIFICATION_LABELS = {
    "TruePositive": "Vrai positif",
    "FalsePositive": "Faux positif",
    "BenignPositive": "Positif bénin",
    "Undetermined": "Indéterminé",
}
CLASSIFICATION_ORDER = ["TruePositive", "FalsePositive", "BenignPositive", "Undetermined"]
CLASSIFICATION_COLORS = {
    "TruePositive": RGBColor(0xB0, 0x20, 0x20),
    "FalsePositive": RGBColor(0xE8, 0x9A, 0x1C),
    "BenignPositive": BRAND_GREEN,
    "Undetermined": RGBColor(0xB7, 0xB7, 0xB7),
}

KPI_LABELS = {
    "MTTA": "Temps moyen de triage",
    "MTTR": "Temps moyen de résolution",
    "MTTC": "Temps moyen de clôture",
}
KPI_ORDER = ["MTTA", "MTTR", "MTTC"]
KPI_ACCENT = {
    "MTTA": RGBColor(0xB0, 0x20, 0x20),
    "MTTR": BRAND_GREEN,
    "MTTC": RGBColor(0x4A, 0x90, 0xC2),
}

# Tailles de police -- augmentees le 22/06/2026 suite au retour visuel
# (titres/legendes/police generale trop petits).
FS_PANEL_HEADER = Pt(24)
FS_PANEL_HEADER_DELTA = Pt(18)
FS_PANEL_SUBHEADER = Pt(19)
FS_PANEL_LINE_LABEL = Pt(20)
FS_PANEL_LINE_VALUE = Pt(20)
FS_PANEL_LINE_DELTA = Pt(18)
FS_DONUT_TITLE = Pt(22)
FS_LEGEND = Pt(17)
FS_DATA_LABEL = Pt(17)
FS_CENTER_TOTAL = Pt(32)
FS_KPI_LABEL = Pt(16)
FS_KPI_VALUE = Pt(32)
FS_KPI_DELTA = Pt(16)

# Le donut visuel (cercle) n'occupe pas toute la largeur de la zone de
# graphique : la legende a droite (chart.legend.position = RIGHT) reserve
# une portion a droite. Pour centrer le TITRE au-dessus du cercle (et non
# au-dessus de cercle+legende), on centre le titre sur cette fraction
# gauche de la largeur plutot que sur la largeur totale. Valeur empirique
# (mesuree sur le rendu reel) -- pas une donnee exposee par python-pptx.
DONUT_VISUAL_WIDTH_RATIO = 0.68


def load_surveillance_history(path: str) -> list:
    """Relit l'onglet Surveillance de l'Excel (cf excel_history.write_surveillance_history)."""
    wb = load_workbook(path, data_only=True)
    ws = wb[SHEET_NAME]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        if row[0] is None:
            continue
        rows.append({
            "Month": str(row[0]),
            "Total": int(row[1] or 0),
            "High": int(row[2] or 0),
            "Medium": int(row[3] or 0),
            "Low": int(row[4] or 0),
            "Informational": int(row[5] or 0),
            "TruePositive": int(row[6] or 0),
            "FalsePositive": int(row[7] or 0),
            "BenignPositive": int(row[8] or 0),
            "Undetermined": int(row[9] or 0),
            "MTTA": float(row[10] or 0.0),
            "MTTR": float(row[11] or 0.0),
            "MTTC": float(row[12] or 0.0),
        })
    return rows


def build_latest_with_evolution(rows: list):
    """
    Isole le DERNIER mois present dans l'historique et le mois precedent
    (s'il existe). Retourne (latest_month, previous_month_ou_None,
    latest_row, previous_row_ou_None).
    """
    if not rows:
        return None, None, None, None
    rows_sorted = sorted(rows, key=lambda r: r["Month"])
    latest = rows_sorted[-1]
    previous = rows_sorted[-2] if len(rows_sorted) >= 2 else None
    return latest["Month"], (previous["Month"] if previous else None), latest, previous


def _format_int_delta(current: int, previous) -> str:
    """Delta signe entre deux compteurs entiers (Total, gravites, categories
    de cloture) : "Nouveau" si pas de mois precedent, sinon "+N"/"-N"/"0"
    -- equivalent entier de _format_duration_delta() pour les durees."""
    if previous is None:
        return "Nouveau"
    delta = current - previous
    if delta > 0:
        return f"+{delta}"
    if delta < 0:
        return str(delta)
    return "0"


def _format_duration(hours: float) -> str:
    """
    Formate une duree exprimee en heures decimales (ex: 2.345) en texte
    lisible pour une carte KPI MTTA/MTTR/MTTC -- decision du 28/06/2026 :
    minutes seules si la duree totale est inferieure a 60 minutes (ex:
    "45 min"), heures + minutes sinon (ex: "2 h 30 min", ou "2 h" si les
    minutes tombent juste). Remplace l'ancien affichage brut en heures
    decimales ("2.345 h").
    """
    total_minutes = round(hours * 60)
    if total_minutes < 60:
        return f"{total_minutes} min"
    h, m = divmod(total_minutes, 60)
    return f"{h} h {m} min" if m else f"{h} h"


def _format_duration_delta(current: float, previous) -> str:
    """Variante signee de _format_duration(), pour l'evolution des cartes KPI.

    Demande du 29/06/2026 : contrairement a _format_duration() (valeur
    courante), cette variante n'affiche JAMAIS les minutes des qu'un ecart
    atteint au moins une heure -- mais ARRONDI a l'heure la plus proche
    (et non une simple troncature) : "-3 h 57 min" devient "-4 h" (57 min
    plus proche de 60 que de 0), tandis que "-3 h 25 min" devient "-3 h"
    (25 min plus proche de 0) -- cf retour utilisateur du 29/06/2026. But :
    garantir un affichage sur EXACTEMENT 2 lignes pour les 3 cartes KPI
    (ligne 1 = libelle, ligne 2 = valeur courante + ecart) ; la
    combinaison valeur courante + ecart complet (h ET min) peut sinon
    deborder sur une 3e ligne pour la carte dont la valeur courante est la
    plus longue (ex: MTTC = "1 h 12 min  -11 h 1 min"), desalignant les 3
    cartes entre elles cote a cote. Applique a TOUTES les cartes (pas
    seulement celle qui debordait) pour rester robuste si l'ecart d'une
    autre carte vient a depasser l'heure un mois donne.
    """
    if previous is None:
        return "Nouveau"
    total_minutes = round((current - previous) * 60)
    if total_minutes == 0:
        return "0 min"
    sign = "+" if total_minutes > 0 else "-"
    abs_minutes = abs(total_minutes)
    if abs_minutes < 60:
        return f"{sign}{abs_minutes} min"
    h, m = divmod(abs_minutes, 60)
    if m >= 30:
        h += 1
    return f"{sign}{h} h"


def _evolution_style(evolution: str) -> dict:
    """Reprend la meme palette que typology_slide (badge vert/rouge/gris)."""
    if evolution == "Nouveau":
        return EVOLUTION_STYLES["new"]
    if evolution in ("0", "0 min"):
        return EVOLUTION_STYLES["flat"]
    if evolution.startswith("+"):
        return EVOLUTION_STYLES["up"]
    return EVOLUTION_STYLES["down"]


def _get_shape_by_name(slide, name):
    """Retourne la shape de la slide portant le nom donne, ou None si aucune
    shape ne correspond (cf typology_slide._get_shape_by_name, duplique ici)."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


# ---------------------------------------------------------------------------
# Panneau de texte (gauche) : titre + repartition gravite + cloture, avec
# badges d'evolution colores (style repris de la slide typologie).
# ---------------------------------------------------------------------------

def _add_panel_line(tf, label: str, value, evolution: str, first: bool = False):
    """Ajoute au panneau de texte une ligne 'label : valeur  (evolution)',
    avec le badge colore de _evolution_style sur l'évolution. `first=True`
    reutilise le tout premier paragraphe du text_frame au lieu d'en
    ajouter un nouveau (deja cree par defaut par python-pptx)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(8)

    run_label = p.add_run()
    run_label.text = f"{label} : "
    run_label.font.name = "Arial"
    run_label.font.size = FS_PANEL_LINE_LABEL
    run_label.font.bold = False
    run_label.font.color.rgb = TEXT_NAVY

    run_value = p.add_run()
    run_value.text = f"{value}  "
    run_value.font.name = "Arial"
    run_value.font.size = FS_PANEL_LINE_VALUE
    run_value.font.bold = True
    run_value.font.color.rgb = TEXT_NAVY

    style = _evolution_style(evolution)
    run_delta = p.add_run()
    run_delta.text = f"({evolution})"
    run_delta.font.name = "Arial"
    run_delta.font.size = FS_PANEL_LINE_DELTA
    run_delta.font.bold = True
    run_delta.font.color.rgb = style["fg"]


def _add_panel_subheader(tf, text: str):
    """Ajoute un sous-titre de section au panneau de texte (ex: 'Répartition
    par gravité', 'Catégorie de clôture')."""
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    p.space_after = Pt(8)
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = FS_PANEL_SUBHEADER
    r.font.bold = True
    r.font.color.rgb = TEXT_NAVY


def _add_panel_spacer(tf):
    """Ajoute un petit paragraphe d'espacement (espace fin) juste après un
    sous-titre de section -- cf _add_panel_section_gap pour l'espacement,
    plus marqué, utilisé avant un sous-titre."""
    p = tf.add_paragraph()
    p.space_after = Pt(2)
    r = p.add_run()
    r.text = " "
    r.font.size = Pt(6)


def _add_panel_section_gap(tf):
    """Espace plus marque avant un sous-titre de section (entre le titre
    et 'Repartition par gravite', et entre la liste gravite et 'Categorie
    de cloture') -- distinct du petit _add_panel_spacer utilise APRES
    chaque sous-titre. Taille calee sur le retour visuel du 22/06/2026."""
    p = tf.add_paragraph()
    p.space_after = Pt(20)
    r = p.add_run()
    r.text = " "
    r.font.size = Pt(8)


def _add_panel_header(tf, text: str, evolution: str = None, first: bool = False):
    """Ajoute l'en-tete principal du panneau de texte (ex: 'N incidents en
    AAAA-MM'), avec un badge d'evolution optionnel accole. `first=True`
    reutilise le tout premier paragraphe du text_frame -- cf _add_panel_line."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(14)
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = FS_PANEL_HEADER
    r.font.bold = True
    r.font.color.rgb = TEXT_NAVY

    if evolution is not None:
        style = _evolution_style(evolution)
        r_delta = p.add_run()
        r_delta.text = f"  ({evolution})"
        r_delta.font.name = "Arial"
        r_delta.font.size = FS_PANEL_HEADER_DELTA
        r_delta.font.bold = True
        r_delta.font.color.rgb = style["fg"]


def _build_text_panel(slide, left, top, width, height, latest: dict, previous):
    """Construit le panneau de texte de gauche (encadré arrondi gris clair) :
    en-tête total + évolution, puis répartition par gravité, puis par
    catégorie de clôture, chaque ligne accompagnée de son badge d'évolution
    vs le mois précédent (cf _add_panel_header / _add_panel_line)."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.adjustments[0] = 0.03
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG
    panel.line.fill.background()
    panel.shadow.inherit = False

    tf = panel.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)

    total_evolution = _format_int_delta(latest["Total"], previous["Total"] if previous else None)
    _add_panel_header(tf, f"{latest['Total']} incidents en {latest['Month']}", evolution=total_evolution, first=True)

    _add_panel_section_gap(tf)
    _add_panel_subheader(tf, "Répartition par gravité")
    _add_panel_spacer(tf)
    for key in SEVERITY_ORDER:
        prev_val = previous[key] if previous else None
        evolution = _format_int_delta(latest[key], prev_val)
        _add_panel_line(tf, SEVERITY_LABELS[key], latest[key], evolution)

    _add_panel_section_gap(tf)
    _add_panel_subheader(tf, "Catégorie de clôture")
    _add_panel_spacer(tf)
    for key in CLASSIFICATION_ORDER:
        prev_val = previous[key] if previous else None
        evolution = _format_int_delta(latest[key], prev_val)
        _add_panel_line(tf, CLASSIFICATION_LABELS[key], latest[key], evolution)

    return panel


# ---------------------------------------------------------------------------
# Donuts (gravite, cloture)
# ---------------------------------------------------------------------------

def _add_donut(slide, left, top, width, height, title: str,
                categories: list, values: list, colors: list, center_total: int):
    """Construit un donut chart natif pptx avec titre centré sur le cercle,
    légende à droite, étiquettes de valeur sur chaque tranche, et un total
    affiché en surimpression au centre de l'anneau (cf commentaires inline
    ci-dessous pour le détail de chaque ajustement visuel)."""
    # Le titre est centre sur la portion VISUELLE du donut (le cercle),
    # pas sur la largeur totale du cadre graphique -- la legende a droite
    # (cf DONUT_VISUAL_WIDTH_RATIO) decale sinon visuellement le titre
    # vers la droite par rapport au cercle.
    title_width = int(width * DONUT_VISUAL_WIDTH_RATIO)
    title_box = slide.shapes.add_textbox(left, top, title_width, Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Arial"
    r.font.size = FS_DONUT_TITLE
    r.font.bold = True
    r.font.color.rgb = TEXT_NAVY

    chart_top = top + Inches(0.5)
    chart_height = height - Inches(0.5)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Incidents", values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, chart_top, width, chart_height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    # NB : include_in_layout=True a ete teste pour reserver explicitement
    # la largeur de la legende (cf retour du 22/06/2026 sur "Informationnelle"
    # qui retournait a la ligne), mais provoque un chevauchement legende/donut
    # au rendu (verifie visuellement) -- revert a False. La taille de police
    # de la legende (FS_LEGEND) a neanmoins ete augmentee, ce qui est la
    # partie surement reproductible de l'ajustement manuel.
    chart.legend.include_in_layout = False
    chart.legend.font.size = FS_LEGEND
    chart.legend.font.name = "Arial"
    chart.has_title = False

    # Layout manuel de la legende : on force explicitement une largeur
    # suffisante pour que "Informationnelle" (la categorie la plus longue)
    # tienne sur une seule ligne -- l'espace par defaut (sans layout
    # manuel) ne s'adapte pas a la taille de police choisie et provoque
    # un retour a la ligne. Fractions exprimees relativement a la zone du
    # graphique (0 a 1), cf retour visuel du 22/06/2026.
    legend_elem = chart.legend._element
    layout_elem = legend_elem.find(qn("c:layout"))
    if layout_elem is None:
        layout_elem = etree.SubElement(legend_elem, qn("c:layout"))
        legend_elem.insert(0, layout_elem)
    manual_layout = etree.SubElement(layout_elem, qn("c:manualLayout"))
    etree.SubElement(manual_layout, qn("c:xMode")).set("val", "edge")
    etree.SubElement(manual_layout, qn("c:yMode")).set("val", "edge")
    etree.SubElement(manual_layout, qn("c:x")).set("val", "0.62")
    etree.SubElement(manual_layout, qn("c:y")).set("val", "0.12")
    etree.SubElement(manual_layout, qn("c:w")).set("val", "0.38")
    etree.SubElement(manual_layout, qn("c:h")).set("val", "0.76")

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    # show_value=True etait manquant : has_data_labels=True seul ne fait que
    # creer l'element <c:dLbls>, il ne suffit pas a activer l'affichage de
    # la valeur (showVal restait a 0) -- c'est pourquoi rien ne s'affichait.
    data_labels.show_value = True
    data_labels.number_format = "0"
    data_labels.number_format_is_linked = False
    data_labels.font.size = FS_DATA_LABEL
    data_labels.font.bold = True
    data_labels.font.color.rgb = TEXT_NAVY

    series = plot.series[0]
    for i, color in enumerate(colors):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color
        point.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        point.format.line.width = Pt(1.5)

    # Largeur de l'anneau (taille du trou central) : valeur Office par defaut ~50.
    plot_elem = plot._element
    hole = plot_elem.find(qn("c:holeSize"))
    if hole is None:
        hole = etree.SubElement(plot_elem, qn("c:holeSize"))
    hole.set("val", "55")

    # Libelle central (ex: total d'incidents) -- overlay statique, non lie
    # nativement aux donnees du donut (python-pptx ne supporte pas les
    # libelles centraux natifs). Centre sur le CERCLE visuel (title_width,
    # meme ratio que le titre ci-dessus), pas sur la largeur totale du
    # cadre (qui inclut la legende) -- sinon le chiffre est decale vers
    # la gauche par rapport au trou du donut.
    center_size = int(min(title_width, chart_height) * 0.42)
    center_left = left + (title_width - center_size) // 2
    center_top = chart_top + (chart_height - center_size) // 2
    center_box = slide.shapes.add_textbox(center_left, center_top, center_size, center_size)
    ctf = center_box.text_frame
    ctf.word_wrap = False
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = 0
    ctf.margin_right = 0
    ctf.margin_top = 0
    ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = str(center_total)
    cr.font.name = "Arial"
    cr.font.size = FS_CENTER_TOTAL
    cr.font.bold = True
    cr.font.color.rgb = TEXT_NAVY


# ---------------------------------------------------------------------------
# Cartes KPI (MTTA / MTTR / MTTC)
# ---------------------------------------------------------------------------

def _add_kpi_card(slide, left, top, width, height, key: str, value: float, previous):
    """Construit une carte KPI (MTTA, MTTR ou MTTC) : encadré blanc à coins
    arrondis avec une bande verticale colorée (KPI_ACCENT), le libellé
    (KPI_LABELS), la valeur formatée (_format_duration) et son évolution
    signée vs le mois précédent (_format_duration_delta)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = KPI_ACCENT[key]
    accent.line.fill.background()
    accent.shadow.inherit = False

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_label = tf.paragraphs[0]
    p_label.alignment = PP_ALIGN.LEFT
    r_label = p_label.add_run()
    r_label.text = KPI_LABELS[key]
    r_label.font.name = "Arial"
    r_label.font.size = FS_KPI_LABEL
    r_label.font.color.rgb = TEXT_MUTED

    p_value = tf.add_paragraph()
    p_value.alignment = PP_ALIGN.LEFT
    r_value = p_value.add_run()
    r_value.text = _format_duration(value)
    r_value.font.name = "Arial"
    r_value.font.size = FS_KPI_VALUE
    r_value.font.bold = True
    r_value.font.color.rgb = TEXT_NAVY

    evolution = _format_duration_delta(value, previous)
    style = _evolution_style(evolution)
    r_delta = p_value.add_run()
    r_delta.text = f"  {evolution}"
    r_delta.font.name = "Arial"
    r_delta.font.size = FS_KPI_DELTA
    r_delta.font.bold = True
    r_delta.font.color.rgb = style["fg"]


# ---------------------------------------------------------------------------
# Assemblage de la slide
# ---------------------------------------------------------------------------

def fill_surveillance_slide(prs, slide, latest_month: str, previous_month, latest: dict, previous):
    """
    Remplit la 3e slide du template (deja presente, titre deja rempli)
    avec : sous-titre (periode), panneau gravite/cloture (gauche), 2 donuts
    (droite haut), 3 cartes KPI MTTA/MTTR/MTTC (droite bas).
    """
    subtitle = _get_shape_by_name(slide, SUBTITLE_SHAPE_NAME)
    if subtitle is not None:
        if previous_month:
            subtitle_text = f"Mois de référence : {latest_month}  (évolution vs {previous_month})"
        else:
            subtitle_text = f"Mois de référence : {latest_month}  (premier mois suivi)"
        # Pas de mise en forme explicite sur le run -- cf typology_slide.
        # fill_evolution_slide pour le detail de la decision du 23/06/2026 :
        # on laisse le run "nu" pour heriter du format Arial Black 33pt
        # navy du placeholder (idx=10 du layout), au lieu du Arial 16pt
        # italique force precedemment (devenu inutile puisqu'on ne reduit
        # plus la taille du sous-titre).
        subtitle.text_frame.text = subtitle_text

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(1.6)

    content_top = Inches(2.35)
    content_bottom = slide_height - Inches(0.6)
    content_height = content_bottom - content_top
    content_width = slide_width - 2 * margin

    left_panel_width = Inches(7.6)
    gap = Inches(0.5)
    right_left = margin + left_panel_width + gap
    right_width = content_width - left_panel_width - gap

    # --- Donuts (droite, haut) ---
    # Taille fixe demandee le 22/06/2026 : 13 x 10 cm (largeur x hauteur),
    # mesuree sur l'objet graphique lui-meme (sans le titre, qui est une
    # zone de texte separee au-dessus). Le bas du cadre reste ancre a la
    # meme position que precedemment (cf kpi_top plus bas, deja calibree
    # pour un bon espacement visuel avec les cartes KPI) -- le cadre etant
    # maintenant plus petit, son HAUT (et donc le titre, juste au-dessus)
    # se retrouve plus bas qu'avant.
    CM = Inches(1 / 2.54)
    chart_height = int(9.5 * CM)
    donut_width = int(13 * CM)
    donut_gap = Inches(0.4)
    title_height = Inches(0.5)

    kpi_top = content_top + Inches(5.65)
    chart_gap_before_kpi = Inches(0.55)
    chart_bottom = kpi_top - chart_gap_before_kpi
    chart_top = chart_bottom - chart_height
    donut_top = chart_top - title_height
    donuts_height = title_height + chart_height

    # --- Panneau de texte (gauche) ---
    # Meme haut que les donuts (donut_top) et meme bas que les cartes KPI
    # (cf kpi_bottom_target) -- cf retour visuel du 22/06/2026 : l'encadre
    # doit matcher la meme hauteur que la colonne de droite.
    kpi_height = Inches(2.2)
    kpi_bottom_target = kpi_top + kpi_height
    left_panel_height = kpi_bottom_target - donut_top
    _build_text_panel(slide, margin, donut_top, left_panel_width, left_panel_height, latest, previous)

    severity_values = [latest[key] for key in SEVERITY_ORDER]
    severity_categories = [SEVERITY_LABELS[key] for key in SEVERITY_ORDER]
    severity_colors = [SEVERITY_COLORS[key] for key in SEVERITY_ORDER]
    _add_donut(
        slide, right_left, donut_top, donut_width, donuts_height,
        "Répartition par gravité", severity_categories, severity_values,
        severity_colors, latest["Total"],
    )

    classification_values = [latest[key] for key in CLASSIFICATION_ORDER]
    classification_categories = [CLASSIFICATION_LABELS[key] for key in CLASSIFICATION_ORDER]
    classification_colors = [CLASSIFICATION_COLORS[key] for key in CLASSIFICATION_ORDER]
    classification_total = sum(classification_values)
    _add_donut(
        slide, right_left + donut_width + donut_gap, donut_top, donut_width, donuts_height,
        "Catégorie de clôture", classification_categories, classification_values,
        classification_colors, classification_total,
    )

    # --- Cartes KPI (droite, bas) ---
    kpi_card_gap = Inches(0.3)
    kpi_card_width = (right_width - 2 * kpi_card_gap) // 3

    for i, key in enumerate(KPI_ORDER):
        kpi_left = right_left + i * (kpi_card_width + kpi_card_gap)
        prev_val = previous[key] if previous else None
        _add_kpi_card(slide, kpi_left, kpi_top, kpi_card_width, kpi_height, key, latest[key], prev_val)

    return slide
