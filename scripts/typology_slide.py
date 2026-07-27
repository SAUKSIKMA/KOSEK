"""
Construction de la slide de synthese "evolution des incidents par typologie",
a partir de l'historique stocke dans l'Excel (cf excel_history.py).

Depuis le 22/06/2026 : la slide n'affiche plus un pivot multi-mois, mais
le DERNIER mois disponible dans l'historique, avec une colonne
"Evolution" comparant chaque typologie au mois precedent (Nouveau / +N /
-N / 0). Le tableau est insere directement dans la 2e slide du template
(template_slide.pptx) -- on NE CREE PAS de nouvelle slide, on reutilise
celle deja presente (titre deja rempli, sous-titre vide a completer).
"""

import copy

from openpyxl import load_workbook
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Couleurs de marque reprises du template (cf placeholders titre/sous-titre) :
# vert "8BC751" et bleu marine "003D56".
BRAND_GREEN = RGBColor(0x8B, 0xC7, 0x51)
BRAND_NAVY = RGBColor(0x00, 0x3D, 0x56)
ROW_BAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_BAND_GREEN_TINT = RGBColor(0xEC, 0xF6, 0xE2)

HEADER_FILL = BRAND_GREEN
HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)

# Regroupement par source d'alerte (slide "Evolution des incidents par
# typologie", decision du 28/06/2026) : bande d'en-tete de groupe et ligne
# de sous-total. Couleur de la bande de groupe changee le 28/06/2026 de
# BRAND_NAVY (#003D56) a #4F81BD, a la demande de l'utilisateur (couleur
# specifique aux "sous-categories" -- titre des sources d'alerte -- pour
# bien la distinguer du navy utilise ailleurs sur les slides de synthese).
GROUP_HEADER_FILL = RGBColor(0x4F, 0x81, 0xBD)
GROUP_HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUBTOTAL_FILL = RGBColor(0xD9, 0xE6, 0xEC)
SUBTOTAL_FONT_COLOR = BRAND_NAVY
GROUP_HEADER_ROW_HEIGHT = Inches(0.45)
SUBTOTAL_ROW_HEIGHT = Inches(0.45)

# Colonne Evolution : couleur de texte + fond "badge" par categorie.
EVOLUTION_STYLES = {
    "up":     {"fg": RGBColor(0xB0, 0x20, 0x20), "bg": RGBColor(0xFB, 0xE1, 0xE1)},  # hausse -> rouge
    "down":   {"fg": RGBColor(0x1F, 0x7A, 0x3D), "bg": RGBColor(0xDC, 0xF0, 0xC9)},  # baisse -> vert
    "new":    {"fg": RGBColor(0x1F, 0x7A, 0x3D), "bg": RGBColor(0xDC, 0xF0, 0xC9)},  # nouveau -> vert
    "flat":   {"fg": RGBColor(0x55, 0x55, 0x55), "bg": RGBColor(0xEC, 0xEC, 0xEC)},  # stable -> gris
}

SUBTITLE_SHAPE_NAME = "Espace réservé du texte 1"

# Nom du shape de TITRE sur les slides de synthese (Evolution, Surveillance,
# SLA) -- decouvert le 28/06/2026 en inspectant le vrai template_slide.pptx :
# le titre n'est PAS un placeholder de type TITLE (slide.shapes.title
# retourne None sur ce template), mais un placeholder de type BODY (idx=11,
# le sous-titre etant idx=10/SUBTITLE_SHAPE_NAME ci-dessus), identifie
# UNIQUEMENT par son nom de shape -- meme principe que SUBTITLE_SHAPE_NAME.
TITLE_SHAPE_NAME = "Espace réservé du texte 2"

# Hauteur de la ligne d'en-tete des tableaux (typologie, SLA) -- decision
# du 23/06/2026 : python-pptx (add_table) distribue par defaut la hauteur
# totale a parts EGALES entre toutes les lignes, ce qui rend la ligne
# d'en-tete (une seule ligne de texte) disproportionnee par rapport aux
# lignes de donnees. On fixe donc explicitement sa hauteur APRES creation
# du tableau ; la hauteur totale du tableau (graphicFrame) se recalcule
# automatiquement comme la somme des hauteurs de ligne -- aucun recalage
# des lignes de donnees n'est necessaire, le tableau se contente d'etre
# moins haut (verifie visuellement, cf retour utilisateur).
HEADER_ROW_HEIGHT = Inches(0.7)

# Pagination multi-slides (decision du 28/06/2026, cf paginate_evolution_
# groups ci-dessous) : python-pptx ne permet pas de connaitre a l'avance
# la hauteur REELLE d'une ligne de tableau une fois le texte rendu (la
# hauteur "h" du XML n'est qu'un MINIMUM, PowerPoint/LibreOffice l'etend
# automatiquement si le texte wrap sur plusieurs lignes -- verifie
# visuellement). On estime donc, UNIQUEMENT a des fins de pagination
# (decider ou couper entre 2 slides), le nombre de lignes qu'occupera un
# intitule de Typologie en fonction de sa longueur, avec une estimation
# volontairement conservative (chars/ligne sous-estime) pour declencher
# un saut de page un peu TROP TOT plutot que TROP TARD -- un peu de blanc
# en bas d'une slide est moins genant qu'un tableau qui debord de la
# slide suivante.
_TYPOLOGY_CHARS_PER_LINE = 55
DATA_ROW_HEIGHT_BASE = Inches(0.38)
DATA_ROW_HEIGHT_PER_EXTRA_LINE = Inches(0.28)


def load_history(path: str) -> list:
    """Relit l'onglet Typologies de l'Excel (cf excel_history.write_history).
    (Onglet renomme de "Historique" a "Typologies" le 22/06/2026, suite a
    l'ajout de l'onglet "Surveillance" dans le meme classeur.)"""
    wb = load_workbook(path, data_only=True)
    ws = wb["Typologies"]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        if row[0] is None:
            continue
        rows.append({
            "Month": str(row[0]),
            "Title": row[1],
            "AlertSources": row[2] or "",
            "IncidentCount": int(row[3] or 0),
        })
    return rows


def build_latest_with_evolution(rows: list):
    """
    Isole le DERNIER mois present dans l'historique et calcule, pour
    chaque typologie de ce mois, son evolution par rapport au mois
    precedent (s'il existe) :
      - "Nouveau" si la typologie n'existait pas le mois precedent
      - "+N" / "-N" / "0" sinon (delta du nombre d'incidents)

    Retourne (latest_month, previous_month_ou_None, table_rows) ou
    table_rows est une liste de dicts {Title, AlertSources, Count,
    Evolution}, triee par nombre d'incidents decroissant.
    """
    months = sorted({r["Month"] for r in rows})
    if not months:
        return None, None, []

    latest_month = months[-1]
    previous_month = months[-2] if len(months) >= 2 else None

    current_rows = [r for r in rows if r["Month"] == latest_month]
    previous_counts = {
        r["Title"]: r["IncidentCount"] for r in rows if r["Month"] == previous_month
    } if previous_month else {}

    table_rows = []
    for r in current_rows:
        title = r["Title"]
        count = r["IncidentCount"]
        if title in previous_counts:
            delta = count - previous_counts[title]
            evolution = f"+{delta}" if delta > 0 else (str(delta) if delta < 0 else "0")
        else:
            evolution = "Nouveau"
        table_rows.append({
            "Title": title,
            "AlertSources": r["AlertSources"],
            "Count": count,
            "Evolution": evolution,
        })

    table_rows.sort(key=lambda r: r["Count"], reverse=True)
    return latest_month, previous_month, table_rows


def group_table_rows_by_source(table_rows: list) -> list:
    """
    Regroupe les lignes de typologie (sortie de build_latest_with_evolution)
    par source d'alerte (champ AlertSources), avec un sous-total (somme
    des Count) par groupe -- demande du 28/06/2026.

    Cas d'une typologie associee a PLUSIEURS sources (chaine deja
    concatenee "A, B" par excel_history._to_display_sources -- rarissime
    en pratique, observe nulle part sur l'historique actuel, mais possible
    par construction puisque typology_normalize.aggregate_typology_rows
    agrege l'UNION des sources par typologie/mois) : la ligne forme son
    propre groupe base sur la chaine combinee exacte, plutot que d'etre
    dupliquee dans chacun des groupes "A" et "B" -- cela eviterait sinon
    un double comptage du nombre d'occurrences dans les sous-totaux
    respectifs.

    table_rows etant deja trie par Count decroissant (cf
    build_latest_with_evolution), l'ordre des typologies au sein de
    chaque groupe est automatiquement deja decroissant (filtrer une
    sequence triee en preserve l'ordre relatif).

    Retourne une liste de dicts {Source, Rows, Subtotal}, triee par
    Subtotal decroissant (source la plus representee en tete).
    """
    groups = {}
    order = []
    for row in table_rows:
        source = row["AlertSources"] or "Source inconnue"
        if source not in groups:
            groups[source] = []
            order.append(source)
        groups[source].append(row)

    result = [
        {"Source": source, "Rows": groups[source], "Subtotal": sum(r["Count"] for r in groups[source])}
        for source in order
    ]
    result.sort(key=lambda g: g["Subtotal"], reverse=True)
    return result


def _estimate_data_row_height(title: str) -> int:
    """Estimation conservative (cf constantes ci-dessus) de la hauteur
    qu'occupera une ligne de typologie, utilisee UNIQUEMENT pour la
    pagination -- pas pour fixer la hauteur reelle de la ligne (laissee
    a l'auto-expansion de PowerPoint/LibreOffice, comme pour le reste du
    tableau)."""
    n_lines = max(1, -(-len(title or "") // _TYPOLOGY_CHARS_PER_LINE))  # division entiere arrondie au sup.
    return DATA_ROW_HEIGHT_BASE + (n_lines - 1) * DATA_ROW_HEIGHT_PER_EXTRA_LINE


def paginate_evolution_groups(groups: list, available_height: int) -> list:
    """
    Decoupe les groupes (sortie de group_table_rows_by_source) en pages,
    chaque page devant tenir dans available_height (EMU) -- decision du
    28/06/2026 : si le tableau groupe ne rentre pas sur une seule slide,
    on continue sur une 2e (ou plus) plutot que de laisser deborder ou de
    tronquer des incidents (chaque ligne reste une typologie nominative,
    on ne masque jamais de donnees pour gagner de la place).

    La ligne d'en-tete de colonnes (Typologie / Nombre d'occurrence /
    Evolution) est re-affichee en haut de CHAQUE page, pour que chacune
    reste lisible independamment des autres.

    Un groupe peut etre scinde entre deux pages si ses lignes ne tiennent
    pas toutes sur la page courante : la bande d'en-tete de groupe est
    alors re-affichee en debut de page suivante, suffixee "(suite)" --
    son sous-total, lui, n'apparait qu'UNE seule fois, juste apres sa
    derniere ligne (donc potentiellement sur la page suivante).

    Retourne une liste de pages, chaque page etant une liste d'items
    types :
      ("group_header", source)
      ("group_header_continued", source)
      ("row", row_dict)
      ("subtotal", source, subtotal_int)
    """
    pages = []
    page = []
    height = HEADER_ROW_HEIGHT  # en-tete de colonnes, reservee sur chaque page

    def start_new_page():
        nonlocal page, height
        pages.append(page)
        page = []
        height = HEADER_ROW_HEIGHT

    for group in groups:
        source = group["Source"]

        if height + GROUP_HEADER_ROW_HEIGHT > available_height and page:
            start_new_page()
        page.append(("group_header", source))
        height += GROUP_HEADER_ROW_HEIGHT

        for row in group["Rows"]:
            row_height = _estimate_data_row_height(row["Title"])
            if height + row_height > available_height and page:
                start_new_page()
                page.append(("group_header_continued", source))
                height += GROUP_HEADER_ROW_HEIGHT
            page.append(("row", row))
            height += row_height

        if height + SUBTOTAL_ROW_HEIGHT > available_height and page:
            start_new_page()
            page.append(("group_header_continued", source))
            height += GROUP_HEADER_ROW_HEIGHT
        page.append(("subtotal", source, group["Subtotal"]))
        height += SUBTOTAL_ROW_HEIGHT

    if page:
        pages.append(page)

    return pages


def clone_slide_from(prs, source_slide):
    """
    Duplique les shapes d'une slide source (titre, sous-titre -- AVANT
    tout tableau ajoute dynamiquement) sur une NOUVELLE slide ajoutee en
    fin de presentation -- utilise pour les slides de continuation
    "(suite)" de la slide d'evolution quand le tableau groupe ne tient
    pas sur une seule slide (cf paginate_evolution_groups, decision du
    28/06/2026).

    A appeler IMPERATIVEMENT avant de remplir la slide source avec son
    propre tableau (sinon le tableau de la source serait lui aussi
    clone). Meme principe que generate_cosec.clone_slide() -- y compris
    la recreation des relations (images, hyperliens...) de la slide
    source sur la nouvelle slide, avec remappage des attributs r:* du XML
    clone (cf bug du 29/06/2026 constate sur generate_cosec.clone_slide :
    sans cela, toute relation referencee par rId dans le spTree copie se
    retrouve invalide -- ou pointant vers la mauvaise cible -- sur la
    nouvelle slide). Cette slide-ci n'a actuellement pas d'image dans son
    spTree (titre/sous-titre uniquement), donc le bug ne s'y manifestait
    pas visuellement, mais le meme correctif est applique par souci de
    robustesse si une image venait a etre ajoutee au template a l'avenir.
    """
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    NOTES_SLIDE_RELTYPE = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    )
    rid_map = {}
    for rel_id, rel in source_slide.part.rels.items():
        if rel.reltype == NOTES_SLIDE_RELTYPE:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel_id] = new_rid

    sp_tree_src = source_slide.shapes._spTree
    sp_tree_dst = new_slide.shapes._spTree

    for el in list(sp_tree_dst):
        sp_tree_dst.remove(el)

    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for el in sp_tree_src:
        new_el = copy.deepcopy(el)
        for sub_el in new_el.iter():
            for attr, value in list(sub_el.attrib.items()):
                if attr.startswith(R_NS) and value in rid_map:
                    sub_el.set(attr, rid_map[value])
        sp_tree_dst.append(new_el)

    return new_slide


def _get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_cell_text(cell, text, size_pt, bold=False, color=None, align=None, bg=None):
    cell.text = str(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    p = cell.text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size_pt)
    run.font.name = "Arial"
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _evolution_style(evolution: str) -> dict:
    if evolution == "Nouveau":
        return EVOLUTION_STYLES["new"]
    if evolution == "0":
        return EVOLUTION_STYLES["flat"]
    if evolution.startswith("+"):
        return EVOLUTION_STYLES["up"]
    return EVOLUTION_STYLES["down"]


def _render_evolution_table(prs, slide, items: list):
    """
    Construit le tableau (3 colonnes : Typologie / Nombre d'occurrence /
    Evolution) pour UNE page donnee, a partir de la liste d'items issue
    de paginate_evolution_groups(). Geometrie identique sur toutes les
    pages (slide principale et continuations) pour un rendu coherent.
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(1.6)

    table_top = Inches(2.6)
    table_width = slide_width - 2 * margin
    table_height = slide_height - table_top - Inches(0.6)

    n_rows = 1 + len(items)
    n_cols = 3

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, margin, table_top, table_width, table_height)
    table = graphic_frame.table

    col_ratios = [0.58, 0.21, 0.21]
    col_widths = [int(table_width * ratio) for ratio in col_ratios]
    col_widths[-1] = table_width - sum(col_widths[:-1])  # absorbe l'arrondi sur la derniere colonne
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Typologie", "Nombre d'occurrence", "Évolution"]
    for c, header in enumerate(headers):
        _set_cell_text(table.cell(0, c), header, 16, bold=True, color=HEADER_FONT_COLOR,
                        align=PP_ALIGN.CENTER, bg=HEADER_FILL)
    table.rows[0].height = HEADER_ROW_HEIGHT

    for r, item in enumerate(items, start=1):
        kind = item[0]

        if kind in ("group_header", "group_header_continued"):
            source = item[1]
            label = source if kind == "group_header" else f"{source} (suite)"
            table.cell(r, 0).merge(table.cell(r, n_cols - 1))
            _set_cell_text(table.cell(r, 0), label, 14, bold=True,
                            color=GROUP_HEADER_FONT_COLOR, align=PP_ALIGN.LEFT, bg=GROUP_HEADER_FILL)
            table.rows[r].height = GROUP_HEADER_ROW_HEIGHT

        elif kind == "row":
            row = item[1]
            # Bandes alternees blanc/vert pale continues sur l'ensemble de
            # la PAGE (pas reinitialisees a chaque groupe ni d'une page a
            # l'autre -- simple alternance par position de ligne r).
            band = ROW_BAND_WHITE if (r % 2 == 1) else ROW_BAND_GREEN_TINT
            _set_cell_text(table.cell(r, 0), row["Title"], 13, bg=band)
            _set_cell_text(table.cell(r, 1), row["Count"], 15, bold=True, align=PP_ALIGN.CENTER, bg=band)

            style = _evolution_style(row["Evolution"])
            _set_cell_text(table.cell(r, 2), row["Evolution"], 17, bold=True,
                            color=style["fg"], align=PP_ALIGN.CENTER, bg=style["bg"])

        elif kind == "subtotal":
            _, _source, subtotal = item
            # Colonnes Nombre d'occurrence + Evolution fusionnees pour
            # centrer la valeur -- le sous-total ne porte que sur le
            # compte d'incidents, pas sur l'evolution.
            table.cell(r, 1).merge(table.cell(r, 2))
            _set_cell_text(table.cell(r, 0), "Sous-total", 13, bold=True,
                            color=SUBTOTAL_FONT_COLOR, align=PP_ALIGN.RIGHT, bg=SUBTOTAL_FILL)
            _set_cell_text(table.cell(r, 1), subtotal, 15, bold=True,
                            color=SUBTOTAL_FONT_COLOR, align=PP_ALIGN.CENTER, bg=SUBTOTAL_FILL)
            table.rows[r].height = SUBTOTAL_ROW_HEIGHT

    return slide


def fill_evolution_slide(prs, evo_slide, latest_month: str, previous_month, table_rows: list):
    """
    Remplit la 2e slide du template (deja presente dans le pptx, PAS
    creee a la volee) avec le tableau Typologie / Nombre d'occurrence /
    Evolution pour le dernier mois disponible, regroupe par source
    d'alerte avec un sous-total par groupe (cf
    group_table_rows_by_source(), demande du 28/06/2026).

    La colonne "Source de l'alerte" (4e colonne jusqu'au 27/06/2026) est
    retiree : elle devient redondante une fois la source affichee en
    bande d'en-tete de chaque groupe, et le tableau passe donc de 4 a 3
    colonnes.

    Pagination automatique (decision du 28/06/2026) : si le tableau
    groupe ne tient pas sur une seule slide, des slides de continuation
    "(suite)" sont clonees a la suite de evo_slide (cf
    paginate_evolution_groups / clone_slide_from) plutot que de laisser
    le tableau deborder ou de tronquer des typologies.

    Le titre de la slide ("Evolution des incidents par typologie") est
    deja renseigne dans le template sur evo_slide -- on ne le touche pas
    (seules les eventuelles slides de continuation voient leur titre
    suffixe "(suite)"). Le sous-titre (vide dans le template) est
    complete avec la periode sur CHAQUE slide produite.

    Retourne la liste de TOUTES les slides "Evolution" produites,
    evo_slide TOUJOURS en premier, suivi des eventuelles slides de
    continuation dans l'ordre. generate_cosec.py doit mettre en avant
    (move_slide_to_front) CHAQUE slide de cette liste, dans l'ordre
    INVERSE, pour preserver leur ordre relatif final (meme principe que
    pour les autres slides de synthese, cf generate_pptx).
    """
    if previous_month:
        base_subtitle_text = f"Mois de référence : {latest_month}  (évolution vs {previous_month})"
    else:
        base_subtitle_text = f"Mois de référence : {latest_month}  (premier mois suivi)"

    # Pas de mise en forme explicite sur le run -- decision du 23/06/2026 :
    # le placeholder herite du layout (lstStyle, idx=10) un format Arial
    # Black 33pt navy concu pour cet usage. On avait initialement (le
    # 22/06/2026) force un Arial 16pt italique en neutralisant le
    # crenage (spc="-360" -> "0") car ce crenage tres negatif, concu
    # pour du 33pt, faisait se chevaucher les lettres a 16pt. Cette
    # neutralisation n'a plus lieu d'etre puisqu'on ne reduit plus la
    # taille : laisser le run "nu" restitue le format du template
    # tel quel, comme c'est deja naturellement le cas pour le sous-titre
    # "Focus sur incident" des slides de detail (jamais touche par le
    # code) -- cf demande utilisateur de reprendre ce meme format pour
    # toutes les slides de synthese.
    subtitle = _get_shape_by_name(evo_slide, SUBTITLE_SHAPE_NAME)
    if subtitle is not None:
        subtitle.text_frame.text = base_subtitle_text

    groups = group_table_rows_by_source(table_rows)

    slide_height = prs.slide_height
    table_top = Inches(2.6)
    available_height = slide_height - table_top - Inches(0.6)

    pages = paginate_evolution_groups(groups, available_height)
    n_pages = len(pages)

    # --- Slides de continuation : clonees AVANT de remplir evo_slide avec
    #     son propre tableau (sinon le tableau serait lui aussi clone). ---
    continuation_slides = [clone_slide_from(prs, evo_slide) for _ in range(n_pages - 1)]

    for i, cont_slide in enumerate(continuation_slides, start=1):
        suite_label = "(suite)" if n_pages <= 2 else f"(suite {i}/{n_pages - 1})"

        title_shape = _get_shape_by_name(cont_slide, TITLE_SHAPE_NAME)
        if title_shape is not None:
            title_shape.text_frame.text = f"{title_shape.text_frame.text} {suite_label}"

        cont_subtitle = _get_shape_by_name(cont_slide, SUBTITLE_SHAPE_NAME)
        if cont_subtitle is not None:
            # Demande du 29/06/2026 : la notion de "(suite)" ne doit
            # apparaitre QUE dans le titre de la slide (suite_label,
            # ci-dessus) -- le sous-titre reste identique a celui de la
            # slide principale, sans suffixe "— suite".
            cont_subtitle.text_frame.text = base_subtitle_text

    all_slides = [evo_slide] + continuation_slides
    for slide, items in zip(all_slides, pages):
        _render_evolution_table(prs, slide, items)

    return all_slides


def move_slide_to_front(prs, slide):
    """Deplace la slide donnee en position 0 dans l'ordre d'affichage du pptx."""
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        if int(sld.get("id")) == slide.slide_id:
            xml_slides.remove(sld)
            xml_slides.insert(0, sld)
            break
